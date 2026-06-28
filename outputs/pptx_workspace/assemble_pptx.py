from pptx import Presentation
from pptx.util import Emu
import os

OUT_DIR = os.path.join(os.path.dirname(__file__), 'slides_browser')
OUT_PPTX = '/home/user/Mahmoud-Gadalla/outputs/Jahizoon_Minister_Presentation.pptx'

prs = Presentation()
prs.slide_width  = Emu(9144000)
prs.slide_height = Emu(5143500)
blank = prs.slide_layouts[6]

for i in range(1, 19):
    num = str(i).zfill(2)
    png = os.path.join(OUT_DIR, f'{num}.png')
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(png, 0, 0, width=prs.slide_width, height=prs.slide_height)
    print(f'  Added slide {num}')

prs.save(OUT_PPTX)
print(f'\n✓ Saved: {OUT_PPTX}')
