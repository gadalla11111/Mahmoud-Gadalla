"""
JAHIZOON — Minister-Ready Presentation
Full matplotlib slide render → PPTX assembly
No HTML, no browser, no RTL bugs.
"""
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as patches
from matplotlib.patches import FancyBboxPatch
import matplotlib.gridspec as gridspec
import numpy as np
from PIL import Image
import io, os, sys

process_env = os.environ.copy()
sys.path.insert(0, '/opt/node22/lib/node_modules')

# ── Fonts ──────────────────────────────────────────────────────────────────────
from matplotlib import font_manager
UNIFONT = "/usr/share/fonts/opentype/unifont/unifont.otf"
font_manager.fontManager.addfont(UNIFONT)
plt.rcParams['font.family'] = 'Unifont'

import arabic_reshaper
from bidi.algorithm import get_display

def ar(t): return get_display(arabic_reshaper.reshape(t))

# ── Brand ──────────────────────────────────────────────────────────────────────
BK  = '#0E0E0E'
GD  = '#C8A24C'
RD  = '#A4232A'
WH  = '#FFFFFF'
DG  = '#1A1A1A'
MG  = '#2A2A2A'
DRD = '#1A0000'
DGRAY = '#111111'

W, H = 13.33, 7.5   # inches at 96dpi → 1280×720 effective
DPI  = 96

OUT_DIR = '/home/user/Mahmoud-Gadalla/outputs/pptx_workspace/slides_final'
os.makedirs(OUT_DIR, exist_ok=True)

def new_slide():
    fig = plt.figure(figsize=(W, H), facecolor=BK)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, W); ax.set_ylim(0, H)
    ax.axis('off')
    ax.set_facecolor(BK)
    return fig, ax

def bar_top(ax, color=GD, h=0.08):
    ax.add_patch(patches.Rectangle((0, H-h), W, h, color=color, zorder=10))

def bar_bottom(ax, color=RD, h=0.06):
    ax.add_patch(patches.Rectangle((0, 0), W, h, color=color, zorder=10))

def bar_left(ax, color=GD, w=0.1):
    ax.add_patch(patches.Rectangle((0, 0), w, H, color=color, zorder=10))

def label(ax, txt, x, y, size=11, color=WH, ha='left', va='center', bold=False, **kw):
    ax.text(x, y, txt, fontsize=size, color=color, ha=ha, va=va,
            fontweight='bold' if bold else 'normal', **kw)

def section_tag(ax, txt, x=0.18, y=H-0.28):
    ax.text(x, y, txt, fontsize=9, color=GD, ha='left', va='center',
            fontproperties=font_manager.FontProperties(family='monospace'),
            letter_spacing=3 if False else None)

def rect(ax, x, y, w, h, color=DGRAY, zorder=1, **kw):
    ax.add_patch(patches.Rectangle((x, y), w, h, color=color, zorder=zorder, **kw))

def divider_h(ax, x, y, w, color=MG, lw=1):
    ax.plot([x, x+w], [y, y], color=color, linewidth=lw, zorder=5)

def divider_v(ax, x, y1, y2, color=MG, lw=1):
    ax.plot([x, x], [y1, y2], color=color, linewidth=lw, zorder=5)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 01 — COVER
# ═══════════════════════════════════════════════════════════════════════════════
def slide01():
    fig, ax = new_slide()
    bar_left(ax, GD, 0.12)
    bar_bottom(ax, RD, 0.05)

    # Right dark panel
    rect(ax, W*0.58, 0, W*0.42, H, '#050505')
    divider_v(ax, W*0.58, 0, H, GD, 2)

    # JAHIZOON — large, vertically centered left
    ax.text(0.22, H*0.60, 'JAHIZOON', fontsize=72, color=GD, ha='left', va='center',
            fontweight='bold', fontfamily='Arial')
    ax.text(0.22, H*0.42, ar('جاهزون'), fontsize=42, color=WH, ha='left', va='center')
    divider_h(ax, 0.22, H*0.32, 4.8, GD, 1.5)
    ax.text(0.22, H*0.24, ar('برنامج التأهيل الوظيفي الوطني للشباب المصري'),
            fontsize=14, color='#888888', ha='left', va='center')

    # Right panel content
    cx = W*0.58 + (W*0.42)/2
    ax.text(cx, H*0.72, 'PRESENTED TO', fontsize=9, color='#444444',
            ha='center', va='center', fontstyle='italic')
    ax.text(cx, H*0.57, ar('وزارة التضامن الاجتماعي'), fontsize=18, color=WH,
            ha='center', va='center', fontweight='bold')
    ax.text(cx, H*0.46, ar('جمهورية مصر العربية'), fontsize=14, color='#888888',
            ha='center', va='center')
    divider_h(ax, W*0.62, H*0.37, W*0.32, RD, 2)
    ax.text(cx, H*0.28, 'June 2026', fontsize=12, color='#555555',
            ha='center', va='center')

    # YM4 label top right
    ax.text(W-0.2, H-0.22, 'YM4 EDUCATION', fontsize=9, color='#333333',
            ha='right', va='top')

    save(fig, '01')

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 02 — THE CRISIS
# ═══════════════════════════════════════════════════════════════════════════════
def slide02():
    fig, ax = new_slide()
    bar_top(ax, GD)
    bar_bottom(ax, RD)

    ax.text(0.5, H-0.3, 'THE CRISIS', fontsize=9, color=GD, ha='left', va='center')
    ax.text(W-0.4, H-0.55, ar('مصر لديها فائض من الشباب · وعجز في التوظيف'),
            fontsize=22, color=WH, ha='right', va='center', fontweight='bold')

    # 3 stat boxes
    boxes = [
        (0.35, '30M',   GD,  ar('شاب مصري في سوق العمل'),  'CAPMAS 2025'),
        (W/3+0.5, '41.5%', RD,  ar('بطالة خريجي الجامعات'),   'CAPMAS Q1 2026'),
        (W*2/3+0.35, '500K', GD,  ar('خريج جديد كل عام'),       ar('وزارة التعليم العالي')),
    ]
    bw = W/3 - 0.5
    for (bx, num, col, lbl, src) in boxes:
        rect(ax, bx, 0.35, bw, H-1.4, '#111111' if col == GD else DRD)
        ax.add_patch(patches.Rectangle((bx, H-1.4), bw, 0.06, color=col))
        ax.text(bx+0.2, H*0.55, num, fontsize=56, color=col, ha='left', va='center',
                fontweight='bold', fontfamily='Arial')
        ax.text(bx+bw/2, H*0.34, lbl, fontsize=13, color=WH, ha='center', va='center')
        ax.text(bx+bw/2, H*0.21, src, fontsize=10, color='#444444', ha='center', va='center')

    ax.text(W-0.4, H*0.08, ar('معدل البطالة الكلية: 6% فقط · الفجوة: ×6.9 مقارنة بالمتوسط العام'),
            fontsize=11, color='#444444', ha='right', va='center')
    save(fig, '02')

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 03 — RESEARCH CONSENSUS
# ═══════════════════════════════════════════════════════════════════════════════
def slide03():
    fig, ax = new_slide()
    bar_top(ax, GD)
    bar_bottom(ax, RD)

    ax.text(0.5, H-0.3, 'RESEARCH CONSENSUS', fontsize=9, color=GD, ha='left')

    # Left — big "13"
    rect(ax, 0.35, 0.35, W*0.38-0.2, H-1.35, '#0A0A0A')
    ax.text(W*0.19, H*0.62, '13', fontsize=120, color=GD, ha='center', va='center',
            fontweight='bold', fontfamily='Arial')
    ax.text(W*0.19, H*0.35, ar('دراسة · نتيجة واحدة'), fontsize=16, color=WH, ha='center', va='center')
    ax.text(W*0.19, H*0.22, 'of 13 studies', fontsize=11, color='#444444', ha='center', va='center')

    # Divider
    divider_v(ax, W*0.38, 0.5, H-0.6, MG, 1.5)

    # Right — findings
    rx = W*0.40
    ax.text(W-0.4, H-0.5,
            ar('"الخريجون يفتقرون إلى المهارات الوظيفية'),
            fontsize=16, color=WH, ha='right', va='center', fontweight='bold')
    ax.text(W-0.4, H-0.78,
            ar('الأساسية التي يحتاجها سوق العمل"'),
            fontsize=16, color=WH, ha='right', va='center', fontweight='bold')

    findings = [
        (GD,  ar('المهارات الناعمة'),    ar('التواصل · العمل الجماعي · إدارة الضغط')),
        (RD,  ar('الفجوة التقنية'),      ar('Excel · العرض · الكتابة المهنية')),
        (GD,  ar('غياب الشبكات'),        ar('80% من الوظائف تُشغَل عبر المعارف')),
    ]
    y = H*0.62
    for col, title, desc in findings:
        ax.add_patch(patches.Rectangle((rx, y-0.1), 0.05, 0.55, color=col))
        ax.text(W-0.4, y+0.3, title, fontsize=14, color=col, ha='right', va='center', fontweight='bold')
        ax.text(W-0.4, y+0.05, desc, fontsize=11, color='#888888', ha='right', va='center')
        y -= 0.78

    ax.text(rx+0.2, H*0.12, 'Consensus.app · 2024–2026', fontsize=10, color='#333333', ha='left')
    save(fig, '03')

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 04 — UNEMPLOYMENT GAP (chart)
# ═══════════════════════════════════════════════════════════════════════════════
def slide04():
    fig = plt.figure(figsize=(W, H), facecolor=BK)

    # Header area
    ax_h = fig.add_axes([0, 0.85, 1, 0.15])
    ax_h.set_facecolor(BK); ax_h.axis('off')
    ax_h.set_xlim(0,1); ax_h.set_ylim(0,1)
    ax_h.add_patch(patches.Rectangle((0, 0.8), 1, 0.2, color=GD))
    ax_h.text(0.03, 0.45, 'THE DATA GAP', fontsize=9, color=GD, ha='left', va='center')
    ax_h.text(0.97, 0.45, ar('فجوة البطالة: الخريجون مقابل سوق العمل'),
              fontsize=20, color=WH, ha='right', va='center', fontweight='bold')

    # Chart area — left 62%
    ax_c = fig.add_axes([0.03, 0.13, 0.57, 0.70], facecolor='#0A0A0A')
    labels = [ar('إجمالي سوق العمل'), ar('خريجو الجامعات')]
    values = [6.0, 41.5]
    colors = [GD, RD]
    bars = ax_c.barh(labels, values, color=colors, height=0.45, edgecolor='none')
    for bar, val, clr in zip(bars, values, ['#0A0A0A', WH]):
        ax_c.text(bar.get_width()-1.2, bar.get_y()+bar.get_height()/2,
                  f'{val}%', va='center', ha='right', fontsize=22,
                  fontweight='bold', color=clr)
    ax_c.set_xlim(0, 55); ax_c.set_xticks([]); ax_c.spines[:].set_visible(False)
    ax_c.tick_params(left=False, labelsize=14)
    ax_c.set_facecolor('#0A0A0A')
    for lbl in ax_c.get_yticklabels(): lbl.set_color(WH)
    ax_c.text(53, -0.65, ar('CAPMAS Q1 2026'), ha='right', fontsize=9, color='#444444', style='italic')

    # Right callout panels — 3 stacked
    panel_data = [
        (DRD, RD,  ar('خريجو الجامعات'), RD,  '41.5%'),
        ('#111111', GD, ar('إجمالي سوق العمل'), GD, '6%'),
        (GD, GD,   ar('الفجوة'),          '#0E0E0E', '×6.9'),
    ]
    heights = [0.28, 0.28, 0.20]
    y0 = 0.13
    for (bg, border, title, numcol, num), ph in zip(panel_data, heights):
        ax_p = fig.add_axes([0.63, y0, 0.34, ph-0.01], facecolor=bg)
        ax_p.axis('off'); ax_p.set_xlim(0,1); ax_p.set_ylim(0,1)
        ax_p.add_patch(patches.Rectangle((0, 0.88), 1, 0.12, color=border))
        ax_p.text(0.95, 0.6, title, fontsize=12, color=WH if bg!=GD else '#0E0E0E',
                  ha='right', va='center')
        ax_p.text(0.08, 0.25, num, fontsize=38, color=numcol if bg!=GD else '#0E0E0E',
                  ha='left', va='center', fontweight='bold', fontfamily='Arial')
        y0 += ph

    # Bottom bar
    ax_b = fig.add_axes([0, 0, 1, 0.04])
    ax_b.set_facecolor(RD); ax_b.axis('off')
    save(fig, '04')

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 05 — EMPLOYER VOICE
# ═══════════════════════════════════════════════════════════════════════════════
def slide05():
    fig, ax = new_slide()
    bar_top(ax, GD)
    bar_bottom(ax, RD)

    ax.text(0.5, H-0.3, 'EMPLOYER VOICE', fontsize=9, color=GD, ha='left')
    ax.text(W-0.4, H-0.52, ar('ماذا يقول أصحاب العمل؟'), fontsize=24, color=WH,
            ha='right', va='center', fontweight='bold')
    ax.text(W-0.4, H-0.82, 'Nexford Employer Survey — Egypt 2026 · N=1,200',
            fontsize=10, color='#444444', ha='right')

    stats = [
        ('78%', RD,  DRD, ar('يرون نقص الكفاءات'), ar('تحدياً رئيسياً للتعيين')),
        ('41%', GD, '#111111', ar('يصفون جاهزية الخريجين'), ar('بـ"دون المستوى"')),
        ('51%', GD, '#111111', ar('مستعدون لتمويل التدريب'), ar('عبر شراكة مؤسسية')),
    ]
    bw = (W - 1.0) / 3
    for i, (num, col, bg, l1, l2) in enumerate(stats):
        bx = 0.35 + i * (bw + 0.15)
        rect(ax, bx, 0.35, bw, H-1.45, bg)
        ax.add_patch(patches.Rectangle((bx, H-1.45), bw, 0.07, color=col))
        ax.text(bx+0.2, H*0.63, num, fontsize=64, color=col, ha='left', va='center',
                fontweight='bold', fontfamily='Arial')
        ax.text(bx+bw/2, H*0.35, l1, fontsize=13, color=WH, ha='center', va='center')
        ax.text(bx+bw/2, H*0.22, l2, fontsize=13, color='#888888', ha='center', va='center')

    save(fig, '05')

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 06 — ROOT CAUSES
# ═══════════════════════════════════════════════════════════════════════════════
def slide06():
    fig, ax = new_slide()
    bar_top(ax, GD)
    bar_bottom(ax, RD)

    ax.text(0.5, H-0.3, 'ROOT CAUSES', fontsize=9, color=GD, ha='left')
    ax.text(W-0.4, H-0.52, ar('لماذا يفشل الخريجون في الحصول على وظائف؟'),
            fontsize=22, color=WH, ha='right', va='center', fontweight='bold')

    causes = [
        ('01', GD, ar('فجوة المهارات الناعمة'),   ar('التواصل · ضبط الوقت · الذكاء العاطفي')),
        ('02', RD, ar('غياب الخبرة العملية'),       ar('مناهج نظرية · لا مشاريع حقيقية')),
        ('03', GD, ar('ضعف إمكانية التحقق'),         ar('أصحاب العمل لا يثقون في الشهادات التقليدية')),
        ('04', RD, ar('انعدام الشبكات المهنية'),     ar('80% من الوظائف تُشغَل عبر المعارف')),
    ]
    cw = (W - 1.0) / 2 - 0.1
    ch = (H - 1.8) / 2 - 0.1
    positions = [(0.35, H*0.15+ch+0.1), (0.35+cw+0.2, H*0.15+ch+0.1),
                 (0.35, H*0.15),          (0.35+cw+0.2, H*0.15)]
    for (cx, cy), (num, col, title, desc) in zip(positions, causes):
        rect(ax, cx, cy, cw, ch, '#111111')
        ax.add_patch(patches.Rectangle((cx, cy+ch-0.06), cw, 0.06, color=col))
        ax.text(cx+0.2, cy+ch*0.72, num, fontsize=28, color=col, ha='left', va='center',
                fontweight='bold', fontfamily='Arial')
        ax.text(cx+cw/2, cy+ch*0.44, title, fontsize=15, color=WH, ha='center', va='center',
                fontweight='bold')
        ax.text(cx+cw/2, cy+ch*0.22, desc, fontsize=11, color='#777777', ha='center', va='center')

    save(fig, '06')

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 07 — TRAJECTORY (chart)
# ═══════════════════════════════════════════════════════════════════════════════
def slide07():
    fig = plt.figure(figsize=(W, H), facecolor=BK)

    ax_h = fig.add_axes([0, 0.85, 1, 0.15])
    ax_h.set_facecolor(BK); ax_h.axis('off')
    ax_h.set_xlim(0,1); ax_h.set_ylim(0,1)
    ax_h.add_patch(patches.Rectangle((0, 0.8), 1, 0.2, color=GD))
    ax_h.text(0.03, 0.45, 'THE TRAJECTORY', fontsize=9, color=GD)
    ax_h.text(0.97, 0.45, ar('المشكلة تتفاقم · كل عام يُضاف نصف مليون'),
              fontsize=20, color=WH, ha='right', fontweight='bold')

    ax_c = fig.add_axes([0.03, 0.12, 0.58, 0.71], facecolor='#0A0A0A')
    years = ['2000','2005','2010','2015','2020','2025','2030','2035']
    vals  = [380, 420, 470, 510, 575, 620, 750, 800]
    bars = ax_c.bar(range(len(years)), vals, color=[GD if v<620 else RD for v in vals],
                    width=0.6, edgecolor='none')
    for bar, v in zip(bars, vals):
        ax_c.text(bar.get_x()+bar.get_width()/2, bar.get_height()+8, str(v),
                  ha='center', va='bottom', fontsize=10, color=WH, fontweight='bold')
    ax_c.set_xticks(range(len(years))); ax_c.set_xticklabels(years, fontsize=11)
    ax_c.set_facecolor('#0A0A0A'); ax_c.set_yticks([])
    ax_c.spines[:].set_visible(False); ax_c.tick_params(bottom=False, colors=WH)

    panels = [
        (DRD, RD,  '800K',  ar('داخل جديد سنوياً بحلول 2035')),
        ('#111111', GD, '+110%', ar('نمو في الضغط على سوق العمل منذ 2000')),
        ('#111111', '#333333', ar('بدون تدخل هيكلي'), ar('ستتضاعف الفجوة')),
    ]
    heights = [0.26, 0.26, 0.19]
    y0 = 0.12
    for (bg, col, num, lbl), ph in zip(panels, heights):
        ax_p = fig.add_axes([0.63, y0, 0.34, ph-0.01], facecolor=bg)
        ax_p.axis('off'); ax_p.set_xlim(0,1); ax_p.set_ylim(0,1)
        ax_p.add_patch(patches.Rectangle((0,0.9),1,0.1,color=col))
        ax_p.text(0.08, 0.5, num, fontsize=32, color=col, ha='left', va='center',
                  fontweight='bold', fontfamily='Arial')
        ax_p.text(0.95, 0.2, lbl, fontsize=11, color=WH, ha='right', va='center')
        y0 += ph

    ax_b = fig.add_axes([0, 0, 1, 0.04]); ax_b.set_facecolor(RD); ax_b.axis('off')
    save(fig, '07')

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 08 — THE ANSWER
# ═══════════════════════════════════════════════════════════════════════════════
def slide08():
    fig, ax = new_slide()
    bar_left(ax, GD, 0.1)
    bar_bottom(ax, RD, 0.05)

    # Left panel
    rect(ax, 0, 0, W*0.42, H, '#040404')
    divider_v(ax, W*0.42, 0, H, '#1A1A1A', 1)

    # Left content — centered
    ax.text(0.22, H*0.80, 'THE ANSWER', fontsize=9, color=GD, ha='left')
    ax.text(0.22, H*0.65, 'JAHIZOON', fontsize=58, color=GD, ha='left', va='center',
            fontweight='bold', fontfamily='Arial')
    ax.text(0.22, H*0.50, ar('جاهزون'), fontsize=34, color=WH, ha='left', va='center')
    divider_h(ax, 0.22, H*0.40, W*0.36, GD, 1.5)
    ax.text(0.22, H*0.30, ar('برنامج وطني يُحوّل الخريجين'), fontsize=13, color='#888888', ha='left')
    ax.text(0.22, H*0.20, ar('إلى كفاءات موثّقة جاهزة للتوظيف الفوري'), fontsize=13, color='#888888', ha='left')

    # Right — 3 pillars
    pillars = [
        (GD, 'TRAIN',  ar('تدريب مكثف · 12 أسبوعاً'), ar('مهارات العمل الفعلية · مدربون معتمدون')),
        (RD, 'ASSESS', ar('تقييم مصوّر · موثّق'),       ar('امتحان نهائي مسجّل يراه أصحاب العمل')),
        (GD, 'PLACE',  ar('توظيف مضمون'),               ar('شبكة شركاء مؤسسيين · ربط مباشر')),
    ]
    ph = (H - 0.1) / 3
    for i, (col, tag, title, desc) in enumerate(pillars):
        y = H - (i+1)*ph
        rect(ax, W*0.44, y, W*0.56-0.2, ph-0.06, '#0D0D0D')
        ax.add_patch(patches.Rectangle((W*0.44, y+ph-0.07), W*0.56-0.2, 0.07, color=col))
        ax.text(W*0.47, y+ph*0.72, tag, fontsize=10, color=col, ha='left', va='center')
        ax.text(W-0.3, y+ph*0.50, title, fontsize=18, color=WH, ha='right', va='center', fontweight='bold')
        ax.text(W-0.3, y+ph*0.25, desc, fontsize=12, color='#666666', ha='right', va='center')
        if i < 2:
            divider_h(ax, W*0.44, y, W*0.56-0.2, '#1A1A1A', 1)

    save(fig, '08')

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 09 — VALUE CHAIN
# ═══════════════════════════════════════════════════════════════════════════════
def slide09():
    fig, ax = new_slide()
    bar_top(ax, GD)
    bar_bottom(ax, RD)

    ax.text(0.5, H-0.3, 'THE PROGRAMME', fontsize=9, color=GD, ha='left')
    ax.text(W-0.4, H-0.52, ar('منهجية جاهزون: من الفجوة إلى التوظيف'),
            fontsize=22, color=WH, ha='right', fontweight='bold')

    steps = [
        ('01', GD,  ar('التقييم'),     ar('تشخيص رقمي')),
        ('02', GD,  ar('التدريب'),     ar('12 أسبوعاً')),
        ('03', RD,  ar('الامتحان'),    ar('تقييم مصوّر')),
        ('04', GD,  ar('التوظيف'),     ar('شركاء مؤسسيون')),
        ('05', GD,  ar('المتابعة'),    ar('6 أشهر دعم')),
    ]
    sw = (W - 1.0) / 5 - 0.1
    sy = 0.35
    sh = H - 1.55

    for i, (num, col, title, desc) in enumerate(steps):
        sx = 0.35 + i*(sw+0.12)
        bg = DRD if col == RD else '#111111'
        rect(ax, sx, sy, sw, sh, bg)
        ax.add_patch(patches.Rectangle((sx, sy+sh-0.07), sw, 0.07, color=col))
        ax.text(sx+sw/2, sy+sh*0.78, num, fontsize=26, color=col, ha='center', va='center',
                fontweight='bold', fontfamily='Arial')
        ax.text(sx+sw/2, sy+sh*0.53, title, fontsize=14, color=WH, ha='center', va='center',
                fontweight='bold')
        ax.text(sx+sw/2, sy+sh*0.30, desc, fontsize=11, color='#888888', ha='center', va='center')
        # Arrow
        if i < 4:
            ax.annotate('', xy=(sx+sw+0.12, sy+sh*0.5), xytext=(sx+sw, sy+sh*0.5),
                        arrowprops=dict(arrowstyle='->', color=GD, lw=2))

    ax.text(W/2, H*0.08, ar('هدف التوظيف: 75% من المتدربين خلال 90 يوماً'),
            fontsize=13, color=GD, ha='center', va='center')
    save(fig, '09')

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — FILMED EXAM
# ═══════════════════════════════════════════════════════════════════════════════
def slide10():
    fig, ax = new_slide()
    bar_left(ax, RD, 0.1)
    bar_bottom(ax, RD, 0.05)

    rect(ax, 0, 0, W*0.42, H, '#040404')
    divider_v(ax, W*0.42, 0, H, '#1A1A1A', 1)

    ax.text(0.22, H*0.82, 'UNIQUE IP', fontsize=9, color=RD, ha='left')
    ax.text(0.22, H*0.68, ar('الامتحان'), fontsize=38, color=WH, ha='left', va='center', fontweight='bold')
    ax.text(0.22, H*0.55, ar('المصوّر'), fontsize=38, color=RD, ha='left', va='center', fontweight='bold')
    divider_h(ax, 0.22, H*0.44, W*0.36, RD, 1.5)
    ax.text(0.22, H*0.34, ar('بدلاً من شهادة ورقية لا يثق بها أحد:'), fontsize=12, color='#888888', ha='left')
    ax.text(0.22, H*0.25, ar('كل متدرب يُقيَّم في موقف عمل حقيقي'), fontsize=12, color='#888888', ha='left')
    ax.text(0.22, H*0.16, ar('ويُسجَّل ويُبثّ لأصحاب العمل مباشرة'), fontsize=12, color='#888888', ha='left')

    # Right — flow + stats
    flow_items = [ar('سيناريو حقيقي'), ar('أداء مسجّل'), ar('توظيف')]
    flow_cols  = [RD, GD, GD]
    fw = W*0.52/3 - 0.15
    for i, (item, col) in enumerate(zip(flow_items, flow_cols)):
        fx = W*0.44 + i*(fw+0.18)
        rect(ax, fx, H*0.62, fw, H*0.24, DRD if col==RD else '#111111')
        ax.add_patch(patches.Rectangle((fx, H*0.62+H*0.24-0.05), fw, 0.05, color=col))
        ax.text(fx+fw/2, H*0.74, item, fontsize=12, color=WH, ha='center', va='center', fontweight='bold')
        if i < 2:
            ax.text(fx+fw+0.09, H*0.74, '→', fontsize=18, color=GD, ha='center', va='center')

    # Stats panel
    rect(ax, W*0.44, 0.35, W*0.55, H*0.24, '#111111')
    pilot = [('45', ar('متدرب')), ('87%', ar('رضا أصحاب العمل')), ('68%', ar('توظيف فعلي'))]
    pw = W*0.55/3
    for i, (num, lbl) in enumerate(pilot):
        px = W*0.44 + i*pw
        ax.text(px+pw/2, H*0.28, num, fontsize=36, color=GD if num!='68%' else RD,
                ha='center', va='center', fontweight='bold', fontfamily='Arial')
        ax.text(px+pw/2, H*0.17, lbl, fontsize=11, color='#888888', ha='center', va='center')
        if i < 2:
            divider_v(ax, px+pw, H*0.12, H*0.38, '#1A1A1A', 1)

    ax.text(W*0.44, H*0.07, ar('أرقام المرحلة التجريبية — نموذج قابل للتوسع'),
            fontsize=10, color='#333333', ha='left')
    save(fig, '10')

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — STAKEHOLDER MATRIX
# ═══════════════════════════════════════════════════════════════════════════════
def slide11():
    fig, ax = new_slide()
    bar_top(ax, GD)
    bar_bottom(ax, RD)

    ax.text(0.5, H-0.3, 'ECOSYSTEM', fontsize=9, color=GD, ha='left')
    ax.text(W-0.4, H-0.52, ar('منظومة الشركاء: ثلاثة محاور · نموذج واحد'),
            fontsize=22, color=WH, ha='right', fontweight='bold')

    cols_data = [
        (GD,  'CORPORATE',  ar('الشركاء المؤسسيون'),
         [ar('تمويل التدريب'), ar('استقبال الخريجين'), ar('تحديد المهارات'), ar('محتوى إعلامي')]),
        (RD,  'GOVERNMENT', ar('الشركاء الحكوميون'),
         [ar('وزارة التضامن'), ar('الجهاز القومي للتشغيل'), ar('مراكز الشباب'), ar('الاعتماد الرسمي')]),
        (GD,  'MEDIA',      ar('شركاء الإعلام'),
         [ar('بث امتحانات النجاح'), ar('قصص تحول حقيقية'), ar('محتوى وطني'), ar('ملايين المشاهدين')]),
    ]
    cw = (W - 1.0) / 3 - 0.1
    for i, (col, tag, title, items) in enumerate(cols_data):
        cx = 0.35 + i*(cw+0.15)
        rect(ax, cx, 0.35, cw, H-1.5, '#111111')
        ax.add_patch(patches.Rectangle((cx, H-1.5), cw, 0.06, color=col))
        ax.text(cx+0.2, H-0.85, tag, fontsize=10, color=col, ha='left', va='center')
        ax.text(cx+cw/2, H-1.18, title, fontsize=15, color=WH, ha='center', va='center', fontweight='bold')
        divider_h(ax, cx+0.1, H-1.56, cw-0.2, '#2A2A2A', 1)
        for j, item in enumerate(items):
            ax.text(cx+cw/2, H*0.50 - j*0.72, item, fontsize=12, color='#AAAAAA',
                    ha='center', va='center')

    ax.text(W/2, H*0.08, ar('كل شريك يربح · جاهزون يُنسّق · الشباب يوظَّف'),
            fontsize=14, color=GD, ha='center', va='center')
    save(fig, '11')

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — BUSINESS MODEL
# ═══════════════════════════════════════════════════════════════════════════════
def slide12():
    fig, ax = new_slide()
    bar_top(ax, GD)
    bar_bottom(ax, RD)

    ax.text(0.5, H-0.3, 'BUSINESS MODEL', fontsize=9, color=GD, ha='left')
    ax.text(W-0.4, H-0.52, ar('نموذج إيرادات متعدد المصادر · مستدام ذاتياً'),
            fontsize=22, color=WH, ha='right', fontweight='bold')

    streams = [
        (GD, '40%', ar('رسوم الشراكة المؤسسية'),
         ar('الشركات تدفع مقابل الوصول إلى الكفاءات المدرّبة')),
        (RD, '35%', ar('إنتاج وتوزيع إعلامي'),
         ar('محتوى قصص النجاح · حقوق البث · الإعلان')),
        (GD, '25%', ar('العقود الحكومية'),
         ar('برامج مموّلة حكومياً · مخرجات موثّقة')),
    ]
    sw = (W*0.62 - 0.5) / 3 - 0.1
    for i, (col, pct, title, desc) in enumerate(streams):
        sx = 0.35 + i*(sw+0.15)
        rect(ax, sx, 0.35, sw, H-1.5, '#111111')
        ax.add_patch(patches.Rectangle((sx, H-1.5), sw, 0.07, color=col))
        ax.text(sx+sw/2, H*0.71, pct, fontsize=44, color=col, ha='center', va='center',
                fontweight='bold', fontfamily='Arial')
        ax.text(sx+sw/2, H*0.50, title, fontsize=13, color=WH, ha='center', va='center', fontweight='bold')
        ax.text(sx+sw/2, H*0.35, desc, fontsize=11, color='#777777', ha='center', va='center')

    # Unit economics box
    rect(ax, W*0.66, 0.35, W*0.32, H-1.5, GD)
    ax.text(W*0.82, H*0.82, ar('اقتصاديات الوحدة'), fontsize=12, color='#0E0E0E',
            ha='center', va='center', fontweight='bold')
    unit = [('3,200', ar('تكلفة المتدرب (ج.م.)')),
            ('5,800', ar('إيراد المتدرب (ج.م.)')),
            ('81%',   ar('هامش الإسهام'))]
    for i, (num, lbl) in enumerate(unit):
        y = H*0.62 - i*0.72
        ax.text(W*0.82, y, num, fontsize=28, color='#0E0E0E', ha='center', va='center',
                fontweight='bold', fontfamily='Arial')
        ax.text(W*0.82, y-0.28, lbl, fontsize=10, color='#3A2A00', ha='center', va='center')

    save(fig, '12')

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — REVENUE MIX (chart placeholder area)
# ═══════════════════════════════════════════════════════════════════════════════
def slide13():
    fig = plt.figure(figsize=(W, H), facecolor=BK)

    ax_h = fig.add_axes([0, 0.85, 1, 0.15])
    ax_h.set_facecolor(BK); ax_h.axis('off')
    ax_h.set_xlim(0,1); ax_h.set_ylim(0,1)
    ax_h.add_patch(patches.Rectangle((0,0.8),1,0.2,color=GD))
    ax_h.text(0.03, 0.45, 'REVENUE MIX', fontsize=9, color=GD)
    ax_h.text(0.97, 0.45, ar('توزيع الإيرادات · محفظة متوازنة'),
              fontsize=20, color=WH, ha='right', fontweight='bold')

    # Donut chart
    ax_c = fig.add_axes([0.03, 0.10, 0.52, 0.73], facecolor=BK)
    vals   = [40, 35, 25]
    colors = [GD, RD, '#555555']
    labels_ar = [ar('شراكات مؤسسية'), ar('إنتاج إعلامي'), ar('عقود حكومية')]
    wedges, _ = ax_c.pie(vals, colors=colors, startangle=90,
                         wedgeprops=dict(width=0.55, edgecolor=BK, linewidth=3))
    ax_c.text(0, 0, ar('الإيرادات'), ha='center', va='center', fontsize=13, color=WH)
    ax_c.set_facecolor(BK)

    # Right legend
    ax_r = fig.add_axes([0.58, 0.10, 0.38, 0.73], facecolor=BK)
    ax_r.axis('off'); ax_r.set_xlim(0,1); ax_r.set_ylim(0,1)
    items = list(zip([40,35,25], colors, labels_ar,
                     [ar('عقود سنوية · أكثر تنبؤاً'),
                      ar('IP فريد · قابل للتوسع'),
                      ar('مضمون · مدعوم بمخرجات')]))
    for i, (pct, col, title, sub) in enumerate(items):
        y = 0.75 - i*0.30
        ax_r.add_patch(patches.Rectangle((0, y+0.04), 0.06, 0.18, color=col))
        ax_r.text(0.95, y+0.18, f'{pct}%', fontsize=32, color=col,
                  ha='right', va='center', fontweight='bold')
        ax_r.text(0.95, y+0.08, title, fontsize=14, color=WH, ha='right', va='center', fontweight='bold')
        ax_r.text(0.95, y-0.02, sub, fontsize=11, color='#666666', ha='right', va='center')

    ax_b = fig.add_axes([0, 0, 1, 0.04]); ax_b.set_facecolor(RD); ax_b.axis('off')
    save(fig, '13')

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — 10-YEAR P&L DASHBOARD (chart)
# ═══════════════════════════════════════════════════════════════════════════════
def slide14():
    YEARS  = ['2025','2026','2027','2028','2029','2030','2031','2032','2033','2034','2035']
    REV    = [2.5,   5.0,   9.2,  15.8,  24.0,  36.5,  52.0,  71.0,  93.0, 115.0, 132.0]
    EBITDA = [-1.2, -0.5,   0.8,   4.2,   8.5,  16.0,  26.0,  38.5,  52.0,  68.0,  68.6]

    fig = plt.figure(figsize=(W, H), facecolor=BK)

    # Header
    ax_h = fig.add_axes([0, 0.85, 1, 0.15])
    ax_h.set_facecolor(BK); ax_h.axis('off'); ax_h.set_xlim(0,1); ax_h.set_ylim(0,1)
    ax_h.add_patch(patches.Rectangle((0,0.8),1,0.2,color=GD))
    ax_h.text(0.03, 0.45, '10-YEAR P&L DASHBOARD', fontsize=9, color=GD)
    ax_h.text(0.97, 0.45, ar('توقعات الأداء المالي 2025–2035 (مليون ج.م.)'),
              fontsize=18, color=WH, ha='right', fontweight='bold')

    # KPI row
    kpis = [
        (GD,  '132M',   ar('الإيراد 2035')),
        (GD,  '68.6M',  ar('EBITDA 2035')),
        (GD,  '30,000', ar('متدرب 2035')),
        (RD,  '2027',   ar('نقطة التعادل')),
    ]
    for i, (col, num, lbl) in enumerate(kpis):
        ax_k = fig.add_axes([0.03+i*0.245, 0.72, 0.235, 0.12],
                            facecolor=DRD if col==RD else '#111111')
        ax_k.axis('off'); ax_k.set_xlim(0,1); ax_k.set_ylim(0,1)
        ax_k.add_patch(patches.Rectangle((0,0.88),1,0.12,color=col))
        ax_k.text(0.05, 0.45, num, fontsize=22, color=col, ha='left', va='center',
                  fontweight='bold', fontfamily='Arial')
        ax_k.text(0.95, 0.2, lbl, fontsize=11, color=WH, ha='right', va='center')

    # Line chart
    ax_c = fig.add_axes([0.03, 0.10, 0.94, 0.60], facecolor='#0A0A0A')
    x = range(len(YEARS))
    ax_c.plot(x, REV,    color=GD, linewidth=3, marker='o', markersize=6, label='Revenue')
    ax_c.plot(x, EBITDA, color=RD, linewidth=3, marker='o', markersize=6, linestyle='--', label='EBITDA')
    ax_c.axhline(0, color=MG, linewidth=1, linestyle=':')
    ax_c.fill_between(x, EBITDA, 0, where=[e>0 for e in EBITDA], alpha=0.15, color=RD)
    ax_c.set_xticks(x); ax_c.set_xticklabels(YEARS, fontsize=10, color=WH)
    ax_c.tick_params(colors=WH, left=True)
    ax_c.yaxis.label.set_color(WH); ax_c.set_facecolor('#0A0A0A')
    for spine in ax_c.spines.values(): spine.set_color(MG)
    ax_c.tick_params(axis='y', labelcolor=WH, labelsize=10)
    ax_c.grid(axis='y', color=MG, linewidth=0.5)
    ax_c.legend(fontsize=11, facecolor='#111111', edgecolor=MG, labelcolor=WH, loc='upper left')

    ax_b = fig.add_axes([0, 0, 1, 0.04]); ax_b.set_facecolor(RD); ax_b.axis('off')
    save(fig, '14')

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — SCALE PLAN
# ═══════════════════════════════════════════════════════════════════════════════
def slide15():
    fig, ax = new_slide()
    bar_top(ax, GD)
    bar_bottom(ax, RD)

    ax.text(0.5, H-0.3, 'SCALE PLAN', fontsize=9, color=GD, ha='left')
    ax.text(W-0.4, H-0.52, ar('من 45 شاباً إلى برنامج وطني'),
            fontsize=24, color=WH, ha='right', fontweight='bold')

    phases = [
        (GD,  'PHASE 01', ar('التجريبية'),  ar('0 – 6 أشهر'),    '45',    ar('متدرب · مركزان'),    ar('3.6 مليون ج.م.')),
        (GD,  'PHASE 02', ar('التوسع'),     ar('6 – 18 شهراً'),   '500',   ar('متدرب · 8 مراكز'),   ar('15 مليون ج.م.')),
        (RD,  'PHASE 03', ar('الوطني'),     ar('18 – 36 شهراً'),  '3,000+',ar('متدرب سنوياً'),       ar('52 مليون ج.م. إيراد')),
    ]
    pw = (W - 1.0) / 3 - 0.2
    for i, (col, tag, title, period, num, unit, rev) in enumerate(phases):
        px = 0.35 + i*(pw+0.3)
        bg = DRD if col==RD else '#111111'
        rect(ax, px, 0.35, pw, H-1.5, bg)
        ax.add_patch(patches.Rectangle((px, H-1.5), pw, 0.07, color=col))
        ax.text(px+pw/2, H*0.84, tag, fontsize=9, color=col, ha='center', va='center')
        ax.text(px+pw/2, H*0.74, title, fontsize=20, color=WH, ha='center', va='center', fontweight='bold')
        ax.text(px+pw/2, H*0.64, period, fontsize=12, color='#888888', ha='center', va='center')
        divider_h(ax, px+0.1, H*0.58, pw-0.2, '#2A2A2A', 1)
        ax.text(px+pw/2, H*0.49, num, fontsize=44, color=col, ha='center', va='center',
                fontweight='bold', fontfamily='Arial')
        ax.text(px+pw/2, H*0.37, unit, fontsize=12, color=WH, ha='center', va='center')
        ax.text(px+pw/2, H*0.19, rev, fontsize=13, color=col, ha='center', va='center', fontweight='bold')
        # Arrow
        if i < 2:
            ax.annotate('', xy=(px+pw+0.3, H*0.5), xytext=(px+pw, H*0.5),
                        arrowprops=dict(arrowstyle='->', color=GD, lw=2.5))

    save(fig, '15')

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — EGYPT VISION 2030
# ═══════════════════════════════════════════════════════════════════════════════
def slide16():
    fig, ax = new_slide()
    bar_top(ax, GD)
    bar_bottom(ax, RD)

    # Background watermark "2030"
    ax.text(W*0.62, H*0.5, '2030', fontsize=180, color='#141414',
            ha='center', va='center', fontweight='bold', fontfamily='Arial', zorder=0)

    ax.text(0.5, H-0.3, 'NATIONAL ALIGNMENT', fontsize=9, color=GD, ha='left')
    ax.text(W-0.4, H-0.52, ar('رؤية مصر 2030 · جاهزون في قلبها'),
            fontsize=22, color=WH, ha='right', fontweight='bold')

    pillars = [
        (GD, '1', ar('التنمية البشرية'),    ar('رفع قابلية التوظيف')),
        (RD, '2', ar('التنويع الاقتصادي'),  ar('قوى عاملة ماهرة')),
        (GD, '3', ar('التحول الرقمي'),      ar('كفاءات رقمية')),
        (GD, '4', ar('التماسك المجتمعي'),   ar('الحد من البطالة')),
    ]
    pw = (W*0.65 - 0.5) / 4 - 0.1
    for i, (col, num, title, sub) in enumerate(pillars):
        px = 0.35 + i*(pw+0.15)
        rect(ax, px, H*0.20, pw, H*0.45, '#111111')
        ax.add_patch(patches.Rectangle((px, H*0.65-0.05), pw, 0.05, color=col))
        ax.text(px+pw/2, H*0.60, num, fontsize=36, color=col, ha='center', va='center',
                fontweight='bold', fontfamily='Arial')
        ax.text(px+pw/2, H*0.46, title, fontsize=13, color=WH, ha='center', va='center', fontweight='bold')
        ax.text(px+pw/2, H*0.34, sub, fontsize=11, color='#777777', ha='center', va='center')

    divider_h(ax, 0.35, H*0.18, W*0.65, MG, 1)
    ax.text(W/2, H*0.12, ar('جاهزون ليس برنامج تدريب · هو استثمار وطني في مستقبل مصر'),
            fontsize=15, color=GD, ha='center', va='center')

    save(fig, '16')

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — THE ASK
# ═══════════════════════════════════════════════════════════════════════════════
def slide17():
    fig, ax = new_slide()
    bar_top(ax, GD)
    bar_bottom(ax, RD)

    ax.text(0.5, H-0.3, 'THE ASK', fontsize=9, color=GD, ha='left')
    ax.text(W-0.4, H-0.52, ar('المرحلة التجريبية · الطلب المحدد'),
            fontsize=22, color=WH, ha='right', fontweight='bold')

    # Budget box - left
    rect(ax, 0.35, H*0.14, W*0.44, H*0.70, GD)
    ax.text(W*0.22+0.35, H*0.75, ar('إجمالي الميزانية'), fontsize=12, color='#0E0E0E',
            ha='center', va='center')
    ax.text(W*0.22+0.35, H*0.62, ar('3,600,000 جنيه'), fontsize=26, color='#0E0E0E',
            ha='center', va='center', fontweight='bold')
    divider_h(ax, 0.45, H*0.57, W*0.42-0.2, '#3A2A00', 1)

    budget = [
        (ar('رواتب المدربين'),         '33.3%'),
        (ar('إيجار وتجهيز المراكز'),   '25.0%'),
        (ar('إدارة المشروع'),           '18.1%'),
        (ar('التقييم المستقل'),          '6.9%'),
        (ar('تطوير المناهج'),           '12.5%'),
        (ar('مستلزمات المتدربين'),      '4.2%'),
    ]
    for i, (item, pct) in enumerate(budget):
        y = H*0.50 - i*0.55
        ax.text(0.5, y, item, fontsize=11, color='#3A2A00', ha='left', va='center')
        ax.text(W*0.44+0.25, y, pct, fontsize=11, color='#0E0E0E', ha='right', va='center', fontweight='bold')

    # Ask items - right
    asks = [
        (GD, '01', ar('تمويل تجريبي'),   ar('3.6 مليون ج.م. لـ 45 متدرباً · 6 أشهر')),
        (GD, '02', ar('وصول للمراكز'),   ar('مراكز الشباب والجهات الحكومية')),
        (RD, '03', ar('اعتماد رسمي'),    ar('شهادة معتمدة تُعزز ثقة أصحاب العمل')),
        (GD, '04', ar('رعاية إعلامية'), ar('دعم الإطلاق وقصص النجاح الأولى')),
    ]
    ah = (H - 1.5) / 4
    for i, (col, num, title, desc) in enumerate(asks):
        ay = H - (i+1)*ah - 0.62
        rect(ax, W*0.47, ay, W*0.50, ah-0.06, '#111111')
        ax.add_patch(patches.Rectangle((W*0.47, ay+ah-0.13), W*0.50, 0.07, color=col))
        ax.text(W*0.49, ay+ah*0.70, num, fontsize=12, color=col, ha='left', va='center', fontweight='bold')
        ax.text(W-0.3, ay+ah*0.52, title, fontsize=14, color=WH, ha='right', va='center', fontweight='bold')
        ax.text(W-0.3, ay+ah*0.25, desc, fontsize=11, color='#777777', ha='right', va='center')

    save(fig, '17')

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — CLOSING
# ═══════════════════════════════════════════════════════════════════════════════
def slide18():
    fig, ax = new_slide()
    bar_left(ax, GD, 0.1)
    bar_bottom(ax, RD, 0.05)
    ax.add_patch(patches.Rectangle((0, H-0.08), W, 0.08, color=GD))

    # Large watermark
    ax.text(W*0.5, H*0.5, ar('جاهزون'), fontsize=140, color='#141414',
            ha='center', va='center', fontweight='bold', zorder=0)

    # Central content
    ax.text(0.22, H*0.75, 'JAHIZOON', fontsize=60, color=GD, ha='left', va='center',
            fontweight='bold', fontfamily='Arial')
    ax.text(0.22, H*0.57, ar('الشباب المصري جاهز للعمل'), fontsize=24, color=WH,
            ha='left', va='center')
    ax.text(0.22, H*0.43, ar('نحتاج فقط من يُعدّهم'), fontsize=24, color='#888888',
            ha='left', va='center')

    divider_h(ax, 0.22, H*0.34, W*0.60, GD, 1.5)
    ax.text(0.22, H*0.24, 'YM4 Education  ·  www.ym4education.com', fontsize=12,
            color=GD, ha='left', va='center')
    ax.text(0.22, H*0.14, 'info@ym4education.com  ·  القاهرة، مصر', fontsize=12,
            color='#666666', ha='left', va='center')

    save(fig, '18')

# ═══════════════════════════════════════════════════════════════════════════════
# ASSEMBLE PPTX
# ═══════════════════════════════════════════════════════════════════════════════
def save(fig, num):
    path = f'{OUT_DIR}/{num}.png'
    fig.savefig(path, dpi=DPI, bbox_inches='tight', facecolor=BK, pad_inches=0)
    plt.close(fig)
    print(f'  Rendered slide {num}')

def assemble_pptx():
    import sys
    sys.path.insert(0, '/opt/node22/lib/node_modules')

    # Use python-pptx
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    import pptx

    prs = Presentation()
    prs.slide_width  = Emu(9144000)   # 10 inches
    prs.slide_height = Emu(5143500)   # 5.625 inches

    blank_layout = prs.slide_layouts[6]  # blank

    for i in range(1, 19):
        num = str(i).zfill(2)
        png = f'{OUT_DIR}/{num}.png'
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(png, 0, 0,
                                 width=prs.slide_width,
                                 height=prs.slide_height)

    out = '/home/user/Mahmoud-Gadalla/outputs/Jahizoon_Minister_Presentation.pptx'
    prs.save(out)
    print(f'\n✓ Saved: {out}')

if __name__ == '__main__':
    print('Rendering slides...')
    slide01(); slide02(); slide03(); slide04(); slide05(); slide06()
    slide07(); slide08(); slide09(); slide10(); slide11(); slide12()
    slide13(); slide14(); slide15(); slide16(); slide17(); slide18()
    print('Assembling PPTX...')
    assemble_pptx()
