"""
MY4 Education — Ministry of Social Solidarity Egypt
Proposal generator: PPTX + DOCX
MERIDIAN brand: #0E0E0E / #FFFFFF / #C8A24C / #A4232A
Arabic RTL throughout — uses real brand assets from guidelines PPTX
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import docx
from docx import Document
from docx.shared import Pt as DPt, RGBColor as DRGBColor, Inches as DInches, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn as dqn
from docx.oxml import OxmlElement
import os

# ── MERIDIAN Colours ──────────────────────────────────────────────────────────
BLACK = RGBColor(0x0E, 0x0E, 0x0E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD  = RGBColor(0xC8, 0xA2, 0x4C)
RED   = RGBColor(0xA4, 0x23, 0x2A)

DBLACK = DRGBColor(0x0E, 0x0E, 0x0E)
DWHITE = DRGBColor(0xFF, 0xFF, 0xFF)
DGOLD  = DRGBColor(0xC8, 0xA2, 0x4C)
DRED   = DRGBColor(0xA4, 0x23, 0x2A)

# ── Slide dimensions ──────────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

# ── Brand asset paths (extracted from brand guidelines PPTX) ──────────────────
ASSETS = "/tmp/pptx_extract/extracted/ppt/media"

def asset(name):
    return os.path.join(ASSETS, name)

# Named asset shortcuts — all 300×780 are woven texture strips (no separate wordmark)
WOVEN_FULL    = asset("image-7-1.png")    # 720×540 — full-slide woven texture
WOVEN_BAR     = asset("image-18-1.png")   # 1320×120 — horizontal woven strip
WOVEN_STRIP   = asset("image-2-1.png")    # 300×780 — vertical woven strip (cover/dividers only)

# Layout constants — full-bleed content, no side strips on content slides
PAD   = Inches(0.55)      # left/right content padding
CONTENT_L = PAD
CONTENT_R = W - PAD
CONTENT_W = CONTENT_R - CONTENT_L

# Icon set (slide 8 — 32 icons, 320×320 RGBA)
ICONS = [asset(f"image-8-{i}.png") for i in range(1, 33)]

# Slide 15 solution icons (5 icons)
SOL_ICONS = [asset(f"image-15-{i}.png") for i in range(1, 6)]

# Slide 10 TOC icons (6 icons)
TOC_ICONS = [asset(f"image-10-{i}.png") for i in range(2, 8)]

# ── Proposal content ──────────────────────────────────────────────────────────
SLIDES = [
    {
        "type": "cover",
        "title": "مبادرة وطنية",
        "subtitle": "لتأهيل الشباب وتمكينهم اقتصادياً",
        "tagline": "من قاعة الدارس إلى سوق العمل",
        "brand": "MY4 Education",
        "ministry": "وزارة التضامن الاجتماعي — جمهورية مصر العربية",
    },
    {
        "type": "toc",
        "title": "فهرس المحتويات",
        "items": [
            ("01", "الإشكالية",        "الشهادة وحدها لا تصنع فرصة عمل"),
            ("02", "الحل والمنهج",     "نموذج مُثبَت ومنهج معتمد"),
            ("03", "إثبات النموذج",    "الأكاديمية العربية — نتائج موثقة"),
            ("04", "الخطة والميزانية", "مسار توسع متدرج بأرقام دقيقة"),
            ("05", "المؤشرات والمتابعة","التزامات قابلة للقياس + بند استرداد"),
            ("06", "MY4 Education",    "الكيان القانوني والفريق والمنافسون"),
            ("07", "الطلب",            "ماذا نطلب من الوزارة بالتحديد"),
        ],
    },
    {
        "type": "exec_summary",
        "section": "ملخص تنفيذي",
        "title": "الأزمة واضحة — الحل موجود — الطلب محدد",
        "body": (
            "تطلب MY4 Education — شركة مصرية مسجلة (س.م.م رقم 157843) — شراكةً رسمية "
            "مع وزارة التضامن الاجتماعي لتشغيل برنامجها المُثبَت في تأهيل الشباب للتوظيف.\n\n"
            "المطلوب: اتفاقية شراكة + إتاحة 3 مراكز تدريبية + إدراج ضمن برنامج ستارت 2026. "
            "التكلفة على الوزارة: 3.6 مليون جنيه للمرحلة التجريبية (300 متدرب / 10 أشهر). "
            "العائد المتوقع: 255 متدرباً موظفاً خلال 90 يوماً من إتمام البرنامج."
        ),
        "stats": [
            ("13.2%", "معدل بطالة الشباب — مصر 2025", "CAPMAS: نشرة سوق العمل Q1 2025"),
            ("16.9%", "بطالة الفئة 20–24 سنة", "CAPMAS: نشرة سوق العمل Q1 2025"),
            ("33.8%", "بطالة الشابات — الأعلى منذ 2015", "CAPMAS: نشرة سوق العمل Q1 2025"),
        ],
    },
    {
        "type": "section_divider",
        "num": "01",
        "title": "الإشكالية",
        "subtitle": "الشهادة وحدها لا تصنع فرصة عمل",
    },
    {
        "type": "problem",
        "section": "01 · الإشكالية",
        "title": "الأرقام تحدد الأزمة بدقة",
        "stats": [
            ("1.2M", "شاب مصري عاطل عن العمل رغم المؤهلات", "CAPMAS 2025"),
            ("16.8%", "بطالة خريجي الجامعات تحديداً", "CAPMAS 2025"),
            ("72%",  "من أصحاب العمل: الخريجون غير مؤهلين عملياً", "CAPMAS 2025"),
            ("31",   "وحدة جامعية تضامن اجتماعي تخدم 250,000 طالب", "وزارة التضامن 2026"),
            ("70%",  "من أسباب الفشل المهني: نقص المهارات لا نقص الفرص", "ILO 2024"),
        ],
    },
    {
        "type": "section_divider",
        "num": "02",
        "title": "الحل والمنهج",
        "subtitle": "نموذج مُثبَت ومنهج معتمد",
    },
    {
        "type": "solution",
        "section": "02 · الحل",
        "title": "كيف يعمل النموذج",
        "steps": [
            ("01", "التشخيص",  "تحليل احتياجات المنطقة وأصحاب العمل — مسار مخصص لكل دفعة قبل بدء التدريب"),
            ("02", "التدريب",  "8 أسابيع مكثفة: مهارات رقمية + ناعمة + محاكاة بيئة العمل الفعلية"),
            ("03", "التوظيف", "ربط مضمون بأصحاب العمل الشركاء + متابعة موثقة 90 يوماً بعد الالتحاق"),
        ],
        "why": (
            "الدمج بين التدريب النظري والتطبيق الفعلي يرفع معدل التوظيف 52%+ "
            "مقارنةً بالتدريب التقليدي (IFC: Workforce Development Report, 2024, p.47). "
            "النموذج مُطبَّق بالفعل مع AASTMT — 200+ متدرب — تقرير التقييم متاح للمراجعة."
        ),
    },
    {
        "type": "curriculum",
        "section": "02 · المنهج",
        "title": "المنهج الدراسي — 8 أسابيع",
        "weeks": [
            ("الأسبوع 1–2", "المهارات الرقمية الأساسية",
             "Microsoft 365 · Google Workspace · إدارة المهام (Trello/Notion) · البريد المهني"),
            ("الأسبوع 3–4", "المهارات الناعمة والتواصل المهني",
             "العروض التقديمية · إدارة الوقت · العمل ضمن فريق · التفاوض والتعامل مع العملاء"),
            ("الأسبوع 5–6", "المسار التخصصي (3 مسارات موازية)",
             "إدارة الأعمال / التسويق الرقمي / دعم تقنية المعلومات — بحسب احتياج سوق العمل المحلي"),
            ("الأسبوع 7–8", "التوظيف والتقييم",
             "إعداد السيرة الذاتية · محاكاة مقابلات العمل · ربط مباشر بأصحاب العمل الشركاء"),
        ],
        "female_track": (
            "المسار النسائي المخصص: تدريب عن بُعد جزئياً + مدربات معتمدات + "
            "جداول مرنة تراعي الالتزامات الأسرية — يستهدف تحديداً الشابات (بطالة 33.8%)"
        ),
        "cert": "اجتيازُ اختبار تقييم مهني معتمد (ECDL-مصر) عند نهاية البرنامج",
    },
    {
        "type": "section_divider",
        "num": "03",
        "title": "إثبات النموذج",
        "subtitle": "الأكاديمية العربية — نتائج موثقة",
    },
    {
        "type": "proof",
        "section": "03 · إثبات النموذج",
        "title": "الأكاديمية العربية — نتائج التجربة الريادية (2024–2025)",
        "body": (
            "شراكة رسمية مع الأكاديمية العربية للعلوم والتكنولوجيا والنقل البحري (AASTMT). "
            "تقرير التقييم المستقل متاح للمراجعة بطلب رسمي. "
            "المقاييس أدناه مستخرجة من تقرير AASTMT لمركز تطوير الكفاءات، مارس 2025."
        ),
        "bullets": [
            "85% من المتدربين وجدوا وظيفة خلال 90 يوماً من إتمام البرنامج (170 من أصل 200)",
            "200+ متدرب في 3 دفعات متتالية — Q2 2024 حتى Q1 2025",
            "4.8 / 5.0 متوسط تقييم رضا المتدربين (استبيان مستقل، n=196)",
            "16 صاحب عمل من القطاعين الخاص والحكومي وقّعوا خطابات توظيف",
        ],
    },
    {
        "type": "employers",
        "section": "03 · شركاء التوظيف",
        "title": "أصحاب العمل الشركاء — الدفعات الثلاث",
        "sectors": [
            ("تقنية المعلومات والاتصالات",
             ["Vodafone Egypt", "Raya Holding", "ITWorx", "Valeo Egypt"]),
            ("الأعمال والإدارة",
             ["Maersk Egypt", "DHL Egypt", "Americana Group", "EGIC"]),
            ("التسويق والإعلام الرقمي",
             ["JWT Egypt", "Procter & Gamble EGY", "Unilever Egypt", "Publicis"]),
            ("القطاع الحكومي والمنظمات",
             ["ITIDA", "MCIT", "Injaz Egypt", "Nahdet El Mahrousa"]),
        ],
        "note": "خطابات النوايا متاحة للمراجعة · يمكن تزويد الوزارة بنسخ رسمية فور طلبها",
    },
    {
        "type": "section_divider",
        "num": "04",
        "title": "الخطة والميزانية",
        "subtitle": "مسار توسع متدرج بأرقام دقيقة",
    },
    {
        "type": "phases",
        "section": "04 · خطة المراحل",
        "title": "مسار التوسع المتدرج",
        "phases": [
            {
                "num": "01", "name": "المرحلة التجريبية",
                "duration": "يناير–أكتوبر ٢٠٢٦", "count": "300 متدرب",
                "points": [
                    "القاهرة · الإسكندرية · الجيزة",
                    "تعيين 12 مدرباً معتمداً (4 لكل مركز)",
                    "تقرير تقييم مستقل عند نهاية المرحلة",
                ]
            },
            {
                "num": "02", "name": "مرحلة التوسع",
                "duration": "يناير–ديسمبر ٢٠٢٧", "count": "1,500 متدرب",
                "points": [
                    "10 مراكز في 6 محافظات",
                    "تفعيل المسار النسائي المخصص",
                    "ربط رسمي ببرنامج ستارت 2026",
                ]
            },
            {
                "num": "03", "name": "الانطلاق الوطني",
                "duration": "٢٠٢٨–٢٠٢٩", "count": "5,000+ متدرب / سنة",
                "points": [
                    "27 محافظة — 31 وحدة جامعية",
                    "شبكة 100+ صاحب عمل",
                    "تمويل ذاتي عبر رسوم التوظيف من الشركات",
                ]
            },
        ],
    },
    {
        "type": "budget",
        "section": "04 · الميزانية",
        "title": "الميزانية التفصيلية — المرحلة الأولى",
        "subtitle": "300 متدرب · 3 مراكز · 10 أشهر",
        "items": [
            ("تطوير المناهج وإعداد المواد التدريبية", "450,000", "12.5%"),
            ("رواتب المدربين (12 مدرباً × 10 أشهر)", "1,200,000", "33.3%"),
            ("إيجار وتجهيز المراكز (3 مراكز × 10 أشهر)", "900,000", "25.0%"),
            ("التقييم المستقل وإعداد التقارير", "250,000", "6.9%"),
            ("مستلزمات ودعم المتدربين (300 × 500 جنيه)", "150,000", "4.2%"),
            ("إدارة المشروع والمصاريف التشغيلية (18%)", "650,000", "18.1%"),
        ],
        "total": "3,600,000",
        "per_trainee": "12,000",
        "roi": "255 متدرباً موظفاً × راتب متوسط 4,500 جنيه/شهر = 1.14M جنيه دخل شهري مُضاف للاقتصاد",
        "clawback": "بند الاسترداد: إذا نسبة التوظيف < 70% — يُعاد 30% من التمويل الحكومي خلال 60 يوماً",
    },
    {
        "type": "section_divider",
        "num": "05",
        "title": "المؤشرات والمتابعة",
        "subtitle": "التزامات قابلة للقياس + بند استرداد",
    },
    {
        "type": "kpis",
        "section": "05 · المؤشرات والمتابعة",
        "title": "التزاماتنا القابلة للقياس",
        "kpis": [
            ("85%+",  "نسبة التوظيف الفعال",    "خلال 90 يوماً من إتمام التدريب"),
            ("4.8/5", "رضا المتدربين",           "تقييم مستقل بعد إتمام البرنامج"),
            ("96%+",  "رضا أصحاب العمل",        "قياس بعد 6 أشهر من التوظيف"),
            ("+52%",  "تحسن معدل التوظيف",      "فوق المتوسط مقارنةً بالتدريب التقليدي"),
        ],
        "qualitative": [
            "لوحة بيانات حية للوزارة تُحدَّث شهرياً (Google Looker Studio)",
            "تقارير ربع سنوية مستقلة من جهة تقييم خارجية معتمدة",
            "بند استرداد رسمي: إذا نسبة التوظيف < 70% → يُعاد 30% من التمويل خلال 60 يوماً",
            "شراكة بحثية مع الجامعات لنشر نتائج النموذج دولياً",
        ],
    },
    {
        "type": "mne",
        "section": "05 · إطار المتابعة",
        "title": "إطار المتابعة والتقييم — M&E Framework",
        "indicators": [
            ("مخرجات", [
                ("عدد المتدربين الملتحقين",    "الهدف: 300 / المرحلة 1",     "سجلات التسجيل — شهرياً"),
                ("معدل إتمام البرنامج",        "الهدف: ≥ 90%",               "نظام الحضور — كل دفعة"),
                ("اجتياز تقييم ECDL",          "الهدف: ≥ 80% من المتدربين",  "شهادات رسمية — كل دفعة"),
            ]),
            ("أثر", [
                ("توظيف خلال 90 يوماً",        "الهدف: ≥ 85%",               "متابعة فردية + خطاب توظيف"),
                ("رضا أصحاب العمل",            "الهدف: ≥ 96%",               "استبيان 6 أشهر بعد التوظيف"),
                ("الاحتفاظ بالوظيفة 6 أشهر",  "الهدف: ≥ 75%",               "اتصال متابعة دوري"),
            ]),
        ],
        "reporting": (
            "التقارير: شهرية للوزارة · ربع سنوية مستقلة · سنوية للنشر العلمي. "
            "المراجعة المستقلة: مركز تقييم معتمد من الوزارة يُحدَّد بالاتفاق."
        ),
    },
    {
        "type": "risk",
        "section": "05 · إدارة المخاطر",
        "title": "تحليل المخاطر وخطة التخفيف",
        "risks": [
            ("ضعف الإقبال في بعض المحافظات",
             "عالية",
             "شراكة مع الوحدات الجامعية للتسويق + حوافز مادية للمتدرب (بدل نقل 200 جنيه/جلسة)"),
            ("تأخر إتاحة المراكز التدريبية",
             "متوسطة",
             "خطة بديلة: عقود إيجار مؤسسات خاصة جاهزة كـ backup — تم تحديدها مسبقاً"),
            ("تذبذب التزام أصحاب العمل",
             "متوسطة",
             "بروتوكولات توظيف موقّعة مسبقاً + قائمة انتظار 30+ شركة بديلة"),
            ("تغيير السياسات الحكومية",
             "منخفضة",
             "هيكل قانوني مرن (س.م.م) + بنود مراجعة سنوية في الاتفاقية"),
            ("نقص المدربين المؤهلين",
             "منخفضة",
             "قاعدة بيانات 45 مدرباً معتمداً — تم بناؤها خلال AASTMT Phase"),
        ],
    },
    {
        "type": "section_divider",
        "num": "06",
        "title": "MY4 Education",
        "subtitle": "الكيان القانوني · الفريق · التميز التنافسي",
    },
    {
        "type": "team",
        "section": "06 · MY4 Education",
        "title": "من نحن — الكيان والفريق",
        "entity": (
            "شركة MY4 Education للتدريب والتطوير المهني (س.م.م رقم 157843) — "
            "مسجلة في مصر، متخصصة في تأهيل الشباب لسوق العمل منذ 2022."
        ),
        "founder": (
            "محمود جاداللا — مؤسس ومدير تنفيذي. "
            "خبرة 8+ سنوات في تصميم برامج التعليم والتدريب المهني. "
            "خريج ومدرب معتمد — شارك في تصميم برامج IFC/World Bank و UNICEF Egypt."
        ),
        "team": [
            ("المدير التنفيذي",          "محمود جاداللا",       "تصميم البرامج + الشراكات الاستراتيجية"),
            ("مدير التدريب",             "د. أميرة حسن",        "مناهج + جودة التدريب — PhD تربية جامعة القاهرة"),
            ("مدير التوظيف وشركاء الأعمال","أحمد الشافعي",    "إدارة شبكة أصحاب العمل — خبرة 10 سنوات HR"),
            ("مدير المتابعة والتقييم",   "سارة مصطفى",          "M&E — شهادة دولية PMD Pro"),
        ],
        "credentials": [
            "شريك معتمد — الأكاديمية العربية للعلوم والتكنولوجيا والنقل البحري (AASTMT)",
            "عضو — منظومة Injaz Egypt / Junior Achievement Worldwide",
            "معتمد — مركز امتحانات ECDL مصر",
        ],
    },
    {
        "type": "competitors",
        "section": "06 · التميز التنافسي",
        "title": "MY4 مقارنةً بالبدائل المتاحة",
        "comparison": [
            ("المعيار",            "MY4 Education",    "ILO/GIZ/UNDP",          "ITIDA/MCIT"),
            ("المدة",              "8 أسابيع",          "6–12 أشهر",             "3–4 أشهر"),
            ("نسبة التوظيف",       "85%+  موثق",        "50–60% تقديري",         "65–70%  بالتقنية فقط"),
            ("التكلفة / متدرب",    "12,000 جنيه",       "25,000–40,000 جنيه",   "18,000–22,000 جنيه"),
            ("المسار النسائي",     "مخصص ومرن",         "محدود",                  "غير متاح"),
            ("بند الاسترداد",      "نعم — 30%",         "لا",                    "لا"),
            ("أثر محلي مُثبَت",    "200+ خريج موثق",    "دراسات فقط",            "نعم — تقنية فقط"),
        ],
        "conclusion": (
            "MY4 Education تقدم الأثر الأعلى بالتكلفة الأدنى مع ضمان مالي مربوط بالنتائج — "
            "لا يوجد مزود آخر في السوق المصرية يجمع هذه العناصر الثلاثة."
        ),
    },
    {
        "type": "section_divider",
        "num": "07",
        "title": "الطلب",
        "subtitle": "ثلاثة بنود محددة تُحقق الشراكة",
    },
    {
        "type": "alignment",
        "section": "07 · الطلب من الوزارة",
        "title": "ماذا نطلب بالتحديد؟",
        "alignment": [
            ("برنامج ستارت ٢٠٢٦",    "شراكة أورنج مصر — MY4 كشريك تنفيذي رسمي للتدريب المهني"),
            ("ملتقى خطوة ٢٠٢٦",      "إدراج MY4 كمزود تدريب معتمد في الملتقى السنوي"),
            ("31 وحدة جامعية تضامن", "إتاحة 3 وحدات (القاهرة/الإسكندرية/الجيزة) لتشغيل المرحلة الأولى"),
            ("رؤية مصر 2030",        "الشراكة تُسهم مباشرةً في مؤشر رفع مشاركة الشباب في سوق العمل"),
        ],
        "requests": [
            {
                "num": "أولاً",
                "title": "اتفاقية شراكة رسمية",
                "body": (
                    "بروتوكول تعاون رسمي بين الوزارة وشركة MY4 Education (س.م.م 157843) "
                    "يُحدد: الصلاحيات · آلية التقييم المستقل · حدود المسؤولية القانونية."
                ),
            },
            {
                "num": "ثانياً",
                "title": "إتاحة 3 مراكز تدريبية",
                "body": (
                    "مساحة تدريبية في الوحدات الجامعية بالقاهرة والإسكندرية والجيزة. "
                    "المطلوب: 2 قاعة × 25 مقعداً في كل مركز · لمدة 10 أشهر."
                ),
            },
            {
                "num": "ثالثاً",
                "title": "الإدراج في برنامج ستارت 2026",
                "body": (
                    "تحويل 300 مستفيد من قاعدة بيانات برنامج ستارت للتدريب لدى MY4 "
                    "مقابل تقديم تقرير التوظيف الرسمي للوزارة خلال 90 يوماً."
                ),
            },
        ],
    },
    {
        "type": "closing",
        "title": "MY4 Education",
        "tagline": "تأهيل حقيقي · توظيف حقيقي · أثر وطني",
        "contact": [
            ("المؤسس",          "محمود جاداللا — Mahmoud Gadalla"),
            ("البريد الإلكتروني", "mahmoud@my4.education"),
            ("الهاتف",          "+20 111 037 111"),
        ],
    },
]

# ═══════════════════════════════════════════════════════════════════════════════
#  UTILITIES
# ═══════════════════════════════════════════════════════════════════════════════


def set_rtl(para_elem):
    pPr = para_elem.find(qn('a:pPr'))
    if pPr is None:
        pPr = etree.SubElement(para_elem, qn('a:pPr'))
        para_elem.insert(0, pPr)
    pPr.set('rtl', '1')
    pPr.set('algn', 'r')


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=None, bg=None,
                 align=PP_ALIGN.RIGHT, rtl=True, font_name="Montserrat",
                 italic=False, line_spacing=None):
    if color is None:
        color = WHITE
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    if rtl:
        set_rtl(p._p)
    if bg:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg
    return txBox


def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=1):
    from pptx.util import Pt as PPt
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = PPt(line_width)
    else:
        shape.line.fill.background()
    return shape


def slide_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_image(slide, path, left, top, width, height=None):
    if height:
        return slide.shapes.add_picture(path, left, top, width, height)
    return slide.shapes.add_picture(path, left, top, width)


def footer(slide, slide_num, total):
    add_rect(slide, 0, H - Inches(0.38), W, Inches(0.38), BLACK)
    add_text_box(slide, "MY4 Education  ·  وزارة التضامن الاجتماعي",
                 Inches(0.4), H - Inches(0.35), Inches(9), Inches(0.32),
                 font_size=7, color=GOLD, align=PP_ALIGN.RIGHT, rtl=True)
    add_text_box(slide, f"{slide_num:02d} / {total:02d}",
                 W - Inches(1.3), H - Inches(0.35), Inches(1.0), Inches(0.32),
                 font_size=7, color=WHITE, align=PP_ALIGN.LEFT, rtl=False, font_name="Inter")


def gold_rule(slide, top, left=None, width_pct=0.88):
    w = W * width_pct
    l = (W - w) / 2 if left is None else left
    add_rect(slide, l, top, w, Pt(2), GOLD)


def section_label(slide, text):
    """Small gold section label top-right."""
    add_text_box(slide, text,
                 Inches(0.4), Inches(0.18), W - Inches(0.8), Inches(0.28),
                 font_size=8, color=GOLD, align=PP_ALIGN.RIGHT, rtl=True, font_name="Inter")


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════════


# ── Shared helpers for content slides ─────────────────────────────────────────

def content_header(slide, section_text, title_text, slide_num, total):
    """Gold top bar, section label, title, gold rule, footer."""
    add_rect(slide, 0, 0, W, Pt(5), GOLD)
    # "M" wordmark top-right
    add_text_box(slide, "M",
                 W - Inches(0.7), Inches(0.04), Inches(0.6), Inches(0.42),
                 font_size=20, bold=True, color=GOLD,
                 align=PP_ALIGN.CENTER, rtl=False, font_name="Montserrat")
    # Section label
    add_text_box(slide, section_text,
                 PAD, Inches(0.1), W - PAD * 2 - Inches(0.7), Inches(0.3),
                 font_size=9, color=GOLD, align=PP_ALIGN.RIGHT,
                 rtl=True, font_name="Inter")
    # Title
    add_text_box(slide, title_text,
                 PAD, Inches(0.42), W - PAD * 2, Inches(0.72),
                 font_size=36, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True)
    # Gold rule under title
    add_rect(slide, PAD, Inches(1.18), CONTENT_W, Pt(2), GOLD)
    footer(slide, slide_num, total)


def slide_cover(slide, data):
    slide_bg(slide, BLACK)
    # Woven texture fills LEFT half (decorative)
    add_image(slide, WOVEN_FULL, 0, 0, W * 0.5, H)
    # Black overlay — left fades into dark
    add_rect(slide, W * 0.3, 0, W * 0.7, H, BLACK)

    # Thin gold top bar full width
    add_rect(slide, 0, 0, W, Pt(5), GOLD)
    # Red left accent
    add_rect(slide, 0, Pt(5), Pt(5), H - Pt(5), RED)

    # Ministry label top-right
    add_text_box(slide, data["ministry"],
                 PAD, Inches(0.2), W - PAD * 2, Inches(0.35),
                 font_size=10, color=GOLD, align=PP_ALIGN.RIGHT, rtl=True, font_name="Inter")

    # "M" mark
    add_text_box(slide, "M",
                 W - Inches(0.9), Inches(0.1), Inches(0.75), Inches(0.55),
                 font_size=24, bold=True, color=GOLD,
                 align=PP_ALIGN.CENTER, rtl=False, font_name="Montserrat")

    # Main title — huge
    add_text_box(slide, data["title"],
                 PAD, Inches(1.5), W - PAD * 2, Inches(1.8),
                 font_size=82, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True)

    add_text_box(slide, data["subtitle"],
                 PAD, Inches(3.35), W - PAD * 2, Inches(0.85),
                 font_size=30, color=GOLD, align=PP_ALIGN.RIGHT, rtl=True, font_name="Inter")

    add_text_box(slide, data["tagline"],
                 PAD, Inches(4.25), W - PAD * 2, Inches(0.5),
                 font_size=16, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True,
                 font_name="Inter", italic=True)

    add_rect(slide, PAD, Inches(4.9), W - PAD * 2, Pt(2), GOLD)

    add_text_box(slide, data["brand"],
                 PAD, Inches(5.05), Inches(5), Inches(0.55),
                 font_size=22, bold=True, color=WHITE, align=PP_ALIGN.RIGHT,
                 rtl=False, font_name="Montserrat")

    add_text_box(slide, "٢٠٢٦",
                 PAD, Inches(5.65), Inches(5), Inches(0.45),
                 font_size=15, color=GOLD, align=PP_ALIGN.RIGHT,
                 rtl=False, font_name="Inter")

    # Narrow woven bar at bottom
    add_image(slide, WOVEN_BAR, 0, H - Inches(0.38), W, Inches(0.38))


def slide_toc(slide, data, slide_num, total):
    slide_bg(slide, BLACK)
    content_header(slide, "", data["title"], slide_num, total)

    row_h = Inches(0.88)
    start_y = Inches(1.28)
    icon_sz = Inches(0.52)

    for i, (num, title, sub) in enumerate(data["items"]):
        y = start_y + i * row_h
        icon_path = TOC_ICONS[i] if i < len(TOC_ICONS) else ICONS[i]

        # Right edge: icon
        add_image(slide, icon_path, CONTENT_R - icon_sz, y + Inches(0.18), icon_sz, icon_sz)

        # Number chip (gold bg, black text)
        add_rect(slide, CONTENT_R - icon_sz - Inches(0.62), y + Inches(0.22),
                 Inches(0.5), Inches(0.38), GOLD)
        add_text_box(slide, num,
                     CONTENT_R - icon_sz - Inches(0.62), y + Inches(0.22),
                     Inches(0.5), Inches(0.38),
                     font_size=11, bold=True, color=BLACK,
                     align=PP_ALIGN.CENTER, rtl=False, font_name="Montserrat")

        text_right = CONTENT_R - icon_sz - Inches(0.75)
        # Title
        add_text_box(slide, title,
                     PAD, y + Inches(0.06), text_right - PAD, Inches(0.4),
                     font_size=16, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True)
        # Subtitle
        add_text_box(slide, sub,
                     PAD, y + Inches(0.46), text_right - PAD, Inches(0.35),
                     font_size=11, color=GOLD, align=PP_ALIGN.RIGHT, rtl=True, font_name="Inter")

        if i < len(data["items"]) - 1:
            add_rect(slide, PAD, y + row_h - Pt(1), CONTENT_W, Pt(1),
                     RGBColor(0x2A, 0x24, 0x14))


def slide_exec_summary(slide, data, slide_num, total):
    slide_bg(slide, BLACK)
    content_header(slide, data["section"], data["title"], slide_num, total)

    add_text_box(slide, data["body"],
                 PAD, Inches(1.32), CONTENT_W * 0.72, Inches(1.9),
                 font_size=13, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True, font_name="Inter")

    # Stats — 3 equal cards
    n = len(data["stats"])
    gap = Inches(0.25)
    card_w = (CONTENT_W - gap * (n - 1)) / n
    card_h = Inches(2.35)
    y = Inches(3.45)

    for i, (num, label, src) in enumerate(data["stats"]):
        x = PAD + i * (card_w + gap)
        add_rect(slide, x, y, card_w, card_h, RGBColor(0x18, 0x14, 0x08),
                 line_color=GOLD, line_width=1)
        # Icon top-left
        add_image(slide, ICONS[i * 4], x + Inches(0.14), y + Inches(0.14),
                  Inches(0.52), Inches(0.52))
        # Big number right-aligned
        add_text_box(slide, num,
                     x + Inches(0.1), y + Inches(0.28), card_w - Inches(0.2), Inches(0.95),
                     font_size=44, bold=True, color=GOLD, align=PP_ALIGN.RIGHT,
                     rtl=False, font_name="Montserrat")
        add_text_box(slide, label,
                     x + Inches(0.1), y + Inches(1.2), card_w - Inches(0.2), Inches(0.7),
                     font_size=12, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True, font_name="Inter")
        add_text_box(slide, src,
                     x + Inches(0.1), y + Inches(1.9), card_w - Inches(0.2), Inches(0.3),
                     font_size=8, color=GOLD, align=PP_ALIGN.RIGHT, rtl=False,
                     font_name="Inter", italic=True)


def slide_section_divider(slide, data, slide_num, total):
    slide_bg(slide, BLACK)
    # Woven fills exactly left half
    add_image(slide, WOVEN_FULL, 0, 0, W * 0.5, H)
    # Clean dark right half
    add_rect(slide, W * 0.5, 0, W * 0.5, H, BLACK)
    # Gold vertical divider line
    add_rect(slide, W * 0.5 - Pt(1.5), 0, Pt(3), H, GOLD)
    # Red top accent bar
    add_rect(slide, 0, 0, W, Pt(5), RED)

    # Ghost section number on left (inside woven)
    add_text_box(slide, data["num"],
                 Inches(0.3), Inches(0.5), W * 0.5 - Inches(0.4), H - Inches(1),
                 font_size=200, bold=True, color=RGBColor(0x20, 0x18, 0x06),
                 align=PP_ALIGN.CENTER, rtl=False, font_name="Montserrat")

    # Section title on right (big, white)
    add_text_box(slide, data["title"],
                 W * 0.52, Inches(2.3), W * 0.44, Inches(1.5),
                 font_size=56, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True)

    # Gold rule
    add_rect(slide, W * 0.52, Inches(3.9), W * 0.44, Pt(3), GOLD)

    # Subtitle
    add_text_box(slide, data["subtitle"],
                 W * 0.52, Inches(4.05), W * 0.44, Inches(0.65),
                 font_size=20, color=GOLD, align=PP_ALIGN.RIGHT, rtl=True, font_name="Inter")

    footer(slide, slide_num, total)


def slide_problem(slide, data, slide_num, total):
    slide_bg(slide, BLACK)
    content_header(slide, data["section"], data["title"], slide_num, total)

    stats = data["stats"]
    # 2-col grid, all cards equal
    gap = Inches(0.22)
    col_w = (CONTENT_W - gap) / 2
    row_h = Inches(1.5)
    start_y = Inches(1.32)

    for i, (num, label, src) in enumerate(stats):
        col = i % 2
        row = i // 2
        x = PAD + col * (col_w + gap)
        y = start_y + row * row_h

        add_rect(slide, x, y, col_w, Inches(1.35),
                 RGBColor(0x1A, 0x14, 0x06), line_color=GOLD, line_width=1)

        icon_idx = (i * 3) % len(ICONS)
        add_image(slide, ICONS[icon_idx], x + Inches(0.12), y + Inches(0.12),
                  Inches(0.5), Inches(0.5))

        add_text_box(slide, num,
                     x + Inches(0.72), y + Inches(0.08), col_w - Inches(0.84), Inches(0.7),
                     font_size=34, bold=True, color=GOLD, align=PP_ALIGN.RIGHT,
                     rtl=False, font_name="Montserrat")

        add_text_box(slide, label,
                     x + Inches(0.12), y + Inches(0.72), col_w - Inches(0.24), Inches(0.48),
                     font_size=11, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True, font_name="Inter")

        add_text_box(slide, src,
                     x + Inches(0.12), y + Inches(1.17), col_w - Inches(0.24), Inches(0.2),
                     font_size=7.5, color=GOLD, align=PP_ALIGN.RIGHT, rtl=False,
                     font_name="Inter", italic=True)


def slide_solution(slide, data, slide_num, total):
    slide_bg(slide, BLACK)
    content_header(slide, data["section"], data["title"], slide_num, total)

    n = len(data["steps"])
    gap = Inches(0.22)
    step_w = (CONTENT_W - gap * (n - 1)) / n
    step_h = Inches(3.6)
    y = Inches(1.32)
    card_colors = [RGBColor(0x1A, 0x14, 0x06),
                   RGBColor(0x20, 0x18, 0x08),
                   RGBColor(0x26, 0x1E, 0x0A)]

    for i, (num, title, body) in enumerate(data["steps"]):
        x = PAD + i * (step_w + gap)
        add_rect(slide, x, y, step_w, step_h, card_colors[i], line_color=GOLD, line_width=1)
        # Gold top strip
        add_rect(slide, x, y, step_w, Inches(0.1), GOLD)

        # Icon centered
        add_image(slide, SOL_ICONS[i], x + step_w / 2 - Inches(0.38), y + Inches(0.2),
                  Inches(0.76), Inches(0.76))

        # Gold number chip
        add_rect(slide, x + Inches(0.14), y + Inches(1.12), Inches(0.52), Inches(0.38), GOLD)
        add_text_box(slide, num,
                     x + Inches(0.14), y + Inches(1.12), Inches(0.52), Inches(0.38),
                     font_size=12, bold=True, color=BLACK,
                     align=PP_ALIGN.CENTER, rtl=False, font_name="Montserrat")

        add_text_box(slide, title,
                     x + Inches(0.14), y + Inches(1.6), step_w - Inches(0.28), Inches(0.55),
                     font_size=17, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True)

        add_text_box(slide, body,
                     x + Inches(0.14), y + Inches(2.22), step_w - Inches(0.28), Inches(1.28),
                     font_size=12, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True, font_name="Inter")

    # Why box
    why_y = Inches(5.08)
    add_rect(slide, PAD, why_y, CONTENT_W, Inches(1.1),
             RGBColor(0x14, 0x10, 0x04), line_color=GOLD, line_width=1)
    add_text_box(slide, data["why"],
                 PAD + Inches(0.18), why_y + Inches(0.1), CONTENT_W - Inches(0.36), Inches(0.9),
                 font_size=11, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True, font_name="Inter")


def slide_proof(slide, data, slide_num, total):
    slide_bg(slide, BLACK)
    content_header(slide, data["section"], data["title"], slide_num, total)

    add_text_box(slide, data["body"],
                 PAD, Inches(1.32), CONTENT_W * 0.72, Inches(1.0),
                 font_size=13, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True, font_name="Inter")

    row_h = Inches(1.0)
    icon_sz = Inches(0.54)

    for i, bullet in enumerate(data["bullets"]):
        y = Inches(2.5) + i * row_h
        icon_idx = (i * 7 + 2) % len(ICONS)

        # Card strip
        add_rect(slide, PAD, y, CONTENT_W, Inches(0.82),
                 RGBColor(0x16, 0x12, 0x05), line_color=RGBColor(0x3A, 0x2E, 0x12))

        # Icon right-aligned inside card
        add_image(slide, ICONS[icon_idx], CONTENT_R - icon_sz - Inches(0.1),
                  y + Inches(0.14), icon_sz, icon_sz)

        # Gold left marker
        add_rect(slide, PAD, y, Pt(4), Inches(0.82), GOLD)

        # Text
        add_text_box(slide, bullet,
                     PAD + Inches(0.18), y + Inches(0.16),
                     CONTENT_W - icon_sz - Inches(0.5), Inches(0.55),
                     font_size=14, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True, font_name="Inter")

    # Woven bar — slim, bottom
    add_image(slide, WOVEN_BAR, PAD, H - Inches(0.72), CONTENT_W, Inches(0.35))


def slide_phases(slide, data, slide_num, total):
    slide_bg(slide, BLACK)
    content_header(slide, data["section"], data["title"], slide_num, total)

    n = len(data["phases"])
    gap = Inches(0.22)
    ph_w = (CONTENT_W - gap * (n - 1)) / n
    ph_h = Inches(5.35)
    y = Inches(1.32)
    accent_colors = [GOLD, RGBColor(0xE0, 0xBC, 0x60), RGBColor(0xA8, 0x82, 0x3C)]

    for i, phase in enumerate(data["phases"]):
        x = PAD + i * (ph_w + gap)
        add_rect(slide, x, y, ph_w, ph_h, RGBColor(0x16, 0x12, 0x06),
                 line_color=accent_colors[i], line_width=2)
        add_rect(slide, x, y, ph_w, Inches(0.12), accent_colors[i])

        # Icon
        add_image(slide, ICONS[i * 6], x + ph_w / 2 - Inches(0.38), y + Inches(0.2),
                  Inches(0.76), Inches(0.76))

        add_text_box(slide, f"المرحلة {phase['num']}",
                     x + Inches(0.12), y + Inches(1.1), ph_w - Inches(0.24), Inches(0.36),
                     font_size=10, color=accent_colors[i], align=PP_ALIGN.RIGHT,
                     rtl=True, font_name="Inter", bold=True)

        add_text_box(slide, phase["name"],
                     x + Inches(0.12), y + Inches(1.48), ph_w - Inches(0.24), Inches(0.52),
                     font_size=15, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True)

        add_rect(slide, x + Inches(0.12), y + Inches(2.06), ph_w - Inches(0.24), Inches(0.32),
                 RGBColor(0x2A, 0x22, 0x0C))
        add_text_box(slide, f"{phase['duration']}  ·  {phase['count']}",
                     x + Inches(0.12), y + Inches(2.06), ph_w - Inches(0.24), Inches(0.32),
                     font_size=9, color=GOLD, align=PP_ALIGN.RIGHT, rtl=True, font_name="Inter")

        add_rect(slide, x + Inches(0.12), y + Inches(2.44), ph_w - Inches(0.24), Pt(1), GOLD)

        for j, pt in enumerate(phase["points"]):
            py = y + Inches(2.6) + j * Inches(0.9)
            add_rect(slide, x + ph_w - Inches(0.28), py + Inches(0.2),
                     Pt(4), Inches(0.24), GOLD)
            add_text_box(slide, pt,
                         x + Inches(0.12), py, ph_w - Inches(0.46), Inches(0.82),
                         font_size=11, color=WHITE, align=PP_ALIGN.RIGHT,
                         rtl=True, font_name="Inter")


def slide_kpis(slide, data, slide_num, total):
    slide_bg(slide, BLACK)
    content_header(slide, data["section"], data["title"], slide_num, total)

    n = len(data["kpis"])
    gap = Inches(0.2)
    kpi_w = (CONTENT_W - gap * (n - 1)) / n
    kpi_h = Inches(2.7)
    y = Inches(1.32)

    for i, (num, label, sub) in enumerate(data["kpis"]):
        x = PAD + i * (kpi_w + gap)
        add_rect(slide, x, y, kpi_w, kpi_h, RGBColor(0x1A, 0x14, 0x06),
                 line_color=GOLD, line_width=1)
        add_image(slide, ICONS[i * 5 + 1], x + kpi_w / 2 - Inches(0.32), y + Inches(0.14),
                  Inches(0.64), Inches(0.64))
        add_text_box(slide, num,
                     x + Inches(0.1), y + Inches(0.88), kpi_w - Inches(0.2), Inches(0.88),
                     font_size=36, bold=True, color=GOLD, align=PP_ALIGN.RIGHT,
                     rtl=False, font_name="Montserrat")
        add_text_box(slide, label,
                     x + Inches(0.1), y + Inches(1.72), kpi_w - Inches(0.2), Inches(0.52),
                     font_size=11, bold=True, color=WHITE, align=PP_ALIGN.RIGHT,
                     rtl=True, font_name="Inter")
        add_text_box(slide, sub,
                     x + Inches(0.1), y + Inches(2.22), kpi_w - Inches(0.2), Inches(0.42),
                     font_size=9, color=GOLD, align=PP_ALIGN.RIGHT, rtl=True, font_name="Inter")

    add_rect(slide, PAD, Inches(4.2), CONTENT_W, Pt(2), GOLD)
    add_text_box(slide, "التزامات إضافية:",
                 PAD, Inches(4.28), CONTENT_W, Inches(0.38),
                 font_size=12, bold=True, color=GOLD, align=PP_ALIGN.RIGHT, rtl=True)

    half_w = (CONTENT_W - Inches(0.2)) / 2
    for i, q in enumerate(data["qualitative"]):
        col = i % 2
        row = i // 2
        x = PAD + col * (half_w + Inches(0.2))
        y = Inches(4.7) + row * Inches(0.55)
        icon_idx = (i * 4 + 3) % len(ICONS)
        add_image(slide, ICONS[icon_idx], x + half_w - Inches(0.52), y,
                  Inches(0.46), Inches(0.46))
        add_text_box(slide, q,
                     x, y + Inches(0.02), half_w - Inches(0.6), Inches(0.46),
                     font_size=11, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True, font_name="Inter")


def slide_alignment(slide, data, slide_num, total):
    slide_bg(slide, BLACK)
    content_header(slide, data["section"], data["title"], slide_num, total)

    gap = Inches(0.2)
    col_w = (CONTENT_W - gap) / 2

    for i, (prog, desc) in enumerate(data["alignment"]):
        col = i % 2
        row = i // 2
        x = PAD + col * (col_w + gap)
        y = Inches(1.34) + row * Inches(1.05)
        add_rect(slide, x, y, col_w, Inches(0.9),
                 RGBColor(0x1A, 0x14, 0x06), line_color=GOLD, line_width=1)
        icon_idx = (i * 8) % len(ICONS)
        add_image(slide, ICONS[icon_idx], x + col_w - Inches(0.58), y + Inches(0.21),
                  Inches(0.48), Inches(0.48))
        add_text_box(slide, prog,
                     x + Inches(0.12), y + Inches(0.07), col_w - Inches(0.74), Inches(0.38),
                     font_size=13, bold=True, color=GOLD, align=PP_ALIGN.RIGHT, rtl=True)
        add_text_box(slide, desc,
                     x + Inches(0.12), y + Inches(0.47), col_w - Inches(0.74), Inches(0.38),
                     font_size=11, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True, font_name="Inter")

    req_y = Inches(3.55)
    add_rect(slide, PAD, req_y, CONTENT_W, Pt(3), RED)
    add_text_box(slide, "المطلوب من الوزارة",
                 PAD, req_y + Inches(0.06), CONTENT_W, Inches(0.42),
                 font_size=15, bold=True, color=RED, align=PP_ALIGN.RIGHT, rtl=True)

    req_w = (CONTENT_W - gap) / 2
    for i, req in enumerate(data["requests"]):
        x = PAD + i * (req_w + gap)
        ry = Inches(4.1)
        add_rect(slide, x, ry, req_w, Inches(2.1), RGBColor(0x20, 0x10, 0x10),
                 line_color=RED, line_width=1)
        add_rect(slide, x, ry, req_w, Inches(0.1), RED)
        add_text_box(slide, req["num"],
                     x + Inches(0.14), ry + Inches(0.14), Inches(0.8), Inches(0.36),
                     font_size=11, bold=True, color=RED, align=PP_ALIGN.RIGHT, rtl=True)
        add_text_box(slide, req["title"],
                     x + Inches(0.14), ry + Inches(0.52), req_w - Inches(0.28), Inches(0.44),
                     font_size=16, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True)
        add_text_box(slide, req["body"],
                     x + Inches(0.14), ry + Inches(1.02), req_w - Inches(0.28), Inches(1.0),
                     font_size=12, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True, font_name="Inter")


def slide_closing(slide, data, slide_num, total):
    slide_bg(slide, BLACK)
    # Woven right half
    add_image(slide, WOVEN_FULL, W * 0.5, 0, W * 0.5, H)
    add_rect(slide, W * 0.5, 0, W * 0.5, H, BLACK)
    # Gold top + left accent
    add_rect(slide, 0, 0, W, Pt(5), GOLD)
    add_rect(slide, 0, Pt(5), Pt(5), H - Pt(5), RED)

    add_text_box(slide, data["title"],
                 PAD, Inches(1.3), W - PAD * 2, Inches(1.8),
                 font_size=82, bold=True, color=WHITE, align=PP_ALIGN.RIGHT,
                 rtl=False, font_name="Montserrat")

    add_rect(slide, PAD, Inches(3.2), W - PAD * 2, Pt(3), GOLD)

    add_text_box(slide, data["tagline"],
                 PAD, Inches(3.32), W - PAD * 2, Inches(0.62),
                 font_size=20, color=GOLD, align=PP_ALIGN.RIGHT, rtl=True, font_name="Inter")

    for i, (label, val) in enumerate(data["contact"]):
        y = Inches(4.18) + i * Inches(0.68)
        icon_idx = 10 + i * 3
        icon_x = CONTENT_R - Inches(0.6)
        add_image(slide, ICONS[icon_idx], icon_x, y + Inches(0.1), Inches(0.5), Inches(0.5))
        add_text_box(slide, f"{label}:  {val}",
                     PAD, y, icon_x - PAD - Inches(0.2), Inches(0.55),
                     font_size=14, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True, font_name="Inter")

    add_image(slide, WOVEN_BAR, 0, H - Inches(0.42), W, Inches(0.42))
    footer(slide, slide_num, total)

def slide_curriculum(slide, data, slide_num, total):
    slide_bg(slide, BLACK)
    content_header(slide, data["section"], data["title"], slide_num, total)

    gap = Inches(0.18)
    row_h = Inches(1.35)
    start_y = Inches(1.32)

    for i, (period, topic, detail) in enumerate(data["weeks"]):
        y = start_y + i * row_h
        card_col = RGBColor(0x1A, 0x14, 0x06) if i % 2 == 0 else RGBColor(0x22, 0x1A, 0x08)
        add_rect(slide, PAD, y, CONTENT_W, Inches(1.22), card_col,
                 line_color=RGBColor(0x3A, 0x2E, 0x12))
        add_rect(slide, PAD, y, Pt(4), Inches(1.22), GOLD)

        # Period chip
        add_rect(slide, PAD + Inches(0.12), y + Inches(0.14), Inches(1.5), Inches(0.34),
                 RGBColor(0x2C, 0x22, 0x0A))
        add_text_box(slide, period,
                     PAD + Inches(0.12), y + Inches(0.14), Inches(1.5), Inches(0.34),
                     font_size=9, bold=True, color=GOLD,
                     align=PP_ALIGN.CENTER, rtl=False, font_name="Inter")

        icon_idx = (i * 6 + 4) % len(ICONS)
        add_image(slide, ICONS[icon_idx], CONTENT_R - Inches(0.62), y + Inches(0.36),
                  Inches(0.5), Inches(0.5))

        add_text_box(slide, topic,
                     PAD + Inches(0.18), y + Inches(0.15), CONTENT_W - Inches(1.3), Inches(0.44),
                     font_size=15, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True)
        add_text_box(slide, detail,
                     PAD + Inches(0.18), y + Inches(0.6), CONTENT_W - Inches(1.3), Inches(0.56),
                     font_size=11, color=GOLD, align=PP_ALIGN.RIGHT, rtl=True, font_name="Inter")

    # Female track note
    note_y = start_y + 4 * row_h + Inches(0.05)
    add_rect(slide, PAD, note_y, CONTENT_W, Inches(0.7),
             RGBColor(0x20, 0x10, 0x10), line_color=RED, line_width=1)
    add_rect(slide, PAD, note_y, Pt(4), Inches(0.7), RED)
    add_text_box(slide, data["female_track"],
                 PAD + Inches(0.18), note_y + Inches(0.08), CONTENT_W - Inches(0.36), Inches(0.56),
                 font_size=11, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True, font_name="Inter")


def slide_employers(slide, data, slide_num, total):
    slide_bg(slide, BLACK)
    content_header(slide, data["section"], data["title"], slide_num, total)

    n = len(data["sectors"])
    gap = Inches(0.2)
    col_w = (CONTENT_W - gap * (n - 1)) / n
    y_start = Inches(1.32)
    col_h = Inches(5.3)

    for i, (sector, companies) in enumerate(data["sectors"]):
        x = PAD + i * (col_w + gap)
        add_rect(slide, x, y_start, col_w, col_h, RGBColor(0x16, 0x12, 0x06),
                 line_color=GOLD, line_width=1)
        add_rect(slide, x, y_start, col_w, Inches(0.08), GOLD)

        icon_idx = (i * 9 + 2) % len(ICONS)
        add_image(slide, ICONS[icon_idx], x + col_w / 2 - Inches(0.3), y_start + Inches(0.12),
                  Inches(0.6), Inches(0.6))

        add_text_box(slide, sector,
                     x + Inches(0.1), y_start + Inches(0.82), col_w - Inches(0.2), Inches(0.6),
                     font_size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER, rtl=True)

        add_rect(slide, x + Inches(0.2), y_start + Inches(1.48), col_w - Inches(0.4), Pt(1), GOLD)

        for j, company in enumerate(companies):
            cy = y_start + Inches(1.62) + j * Inches(0.88)
            add_rect(slide, x + Inches(0.12), cy, col_w - Inches(0.24), Inches(0.72),
                     RGBColor(0x22, 0x1C, 0x0A))
            add_text_box(slide, company,
                         x + Inches(0.12), cy + Inches(0.18), col_w - Inches(0.24), Inches(0.44),
                         font_size=12, color=WHITE, align=PP_ALIGN.CENTER, rtl=False, font_name="Inter")

    add_text_box(slide, data["note"],
                 PAD, H - Inches(0.95), CONTENT_W, Inches(0.3),
                 font_size=9, color=GOLD, align=PP_ALIGN.RIGHT, rtl=True,
                 font_name="Inter", italic=True)


def slide_budget(slide, data, slide_num, total):
    slide_bg(slide, BLACK)
    content_header(slide, data["section"], data["title"], slide_num, total)

    add_text_box(slide, data["subtitle"],
                 PAD, Inches(1.22), CONTENT_W, Inches(0.26),
                 font_size=11, color=GOLD, align=PP_ALIGN.RIGHT, rtl=True, font_name="Inter")

    # Table header
    add_rect(slide, PAD, Inches(1.54), CONTENT_W, Inches(0.36), GOLD)
    header_cols = [("البند", 0.55), ("المبلغ (جنيه)", 0.22), ("النسبة", 0.23)]
    x_cur = PAD
    for label, frac in header_cols:
        cw = CONTENT_W * frac
        add_text_box(slide, label, x_cur, Inches(1.54), cw, Inches(0.36),
                     font_size=11, bold=True, color=BLACK,
                     align=PP_ALIGN.CENTER, rtl=True, font_name="Montserrat")
        x_cur += cw

    row_h = Inches(0.6)
    for i, (item, amount, pct) in enumerate(data["items"]):
        y = Inches(1.9) + i * row_h
        row_col = RGBColor(0x1A, 0x14, 0x06) if i % 2 == 0 else RGBColor(0x22, 0x1A, 0x08)
        add_rect(slide, PAD, y, CONTENT_W, row_h - Pt(1), row_col)
        x_cur = PAD
        for val, frac in [(item, 0.55), (amount, 0.22), (pct, 0.23)]:
            cw = CONTENT_W * frac
            add_text_box(slide, val, x_cur, y + Inches(0.1), cw, row_h - Inches(0.14),
                         font_size=11, color=WHITE, align=PP_ALIGN.CENTER,
                         rtl=(frac == 0.55), font_name="Inter")
            x_cur += cw

    # Total row
    tot_y = Inches(1.9) + len(data["items"]) * row_h
    add_rect(slide, PAD, tot_y, CONTENT_W, Inches(0.48), RGBColor(0x2C, 0x22, 0x0A),
             line_color=GOLD, line_width=2)
    add_text_box(slide, f"الإجمالي:  {data['total']} جنيه  |  {data['per_trainee']} جنيه / متدرب",
                 PAD, tot_y + Inches(0.08), CONTENT_W, Inches(0.34),
                 font_size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER, rtl=True, font_name="Montserrat")

    # ROI + clawback
    add_text_box(slide, data["roi"],
                 PAD, tot_y + Inches(0.58), CONTENT_W, Inches(0.32),
                 font_size=10, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True, font_name="Inter")
    add_rect(slide, PAD, tot_y + Inches(0.96), CONTENT_W, Inches(0.42),
             RGBColor(0x20, 0x10, 0x10), line_color=RED, line_width=1)
    add_text_box(slide, data["clawback"],
                 PAD + Inches(0.12), tot_y + Inches(1.02), CONTENT_W - Inches(0.24), Inches(0.32),
                 font_size=10, bold=True, color=RED, align=PP_ALIGN.RIGHT, rtl=True, font_name="Inter")


def slide_mne(slide, data, slide_num, total):
    slide_bg(slide, BLACK)
    content_header(slide, data["section"], data["title"], slide_num, total)

    y = Inches(1.32)
    half_w = (CONTENT_W - Inches(0.25)) / 2

    for col_idx, (cat, rows) in enumerate(data["indicators"]):
        x = PAD + col_idx * (half_w + Inches(0.25))
        add_rect(slide, x, y, half_w, Inches(0.36), GOLD)
        add_text_box(slide, cat, x, y, half_w, Inches(0.36),
                     font_size=12, bold=True, color=BLACK,
                     align=PP_ALIGN.CENTER, rtl=True, font_name="Montserrat")

        for j, (indicator, target, method) in enumerate(rows):
            ry = y + Inches(0.4) + j * Inches(1.02)
            row_col = RGBColor(0x1A, 0x14, 0x06) if j % 2 == 0 else RGBColor(0x22, 0x1A, 0x08)
            add_rect(slide, x, ry, half_w, Inches(0.96), row_col,
                     line_color=RGBColor(0x3A, 0x2E, 0x12))
            add_text_box(slide, indicator, x + Inches(0.1), ry + Inches(0.06),
                         half_w - Inches(0.2), Inches(0.38),
                         font_size=12, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True)
            add_text_box(slide, target, x + Inches(0.1), ry + Inches(0.44),
                         half_w - Inches(0.2), Inches(0.28),
                         font_size=10, color=GOLD, align=PP_ALIGN.RIGHT, rtl=True, font_name="Inter")
            add_text_box(slide, method, x + Inches(0.1), ry + Inches(0.68),
                         half_w - Inches(0.2), Inches(0.24),
                         font_size=9, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True,
                         font_name="Inter", italic=True)

    rep_y = Inches(4.62)
    add_rect(slide, PAD, rep_y, CONTENT_W, Inches(0.74),
             RGBColor(0x14, 0x10, 0x04), line_color=GOLD, line_width=1)
    add_text_box(slide, data["reporting"],
                 PAD + Inches(0.18), rep_y + Inches(0.1), CONTENT_W - Inches(0.36), Inches(0.58),
                 font_size=11, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True, font_name="Inter")


def slide_risk(slide, data, slide_num, total):
    slide_bg(slide, BLACK)
    content_header(slide, data["section"], data["title"], slide_num, total)

    RISK_COLORS = {
        "عالية":    RGBColor(0xA4, 0x23, 0x2A),
        "متوسطة":   RGBColor(0xC8, 0x80, 0x14),
        "منخفضة":   RGBColor(0x28, 0x82, 0x40),
    }

    row_h = Inches(1.06)
    for i, (risk, level, mitigation) in enumerate(data["risks"]):
        y = Inches(1.32) + i * row_h
        add_rect(slide, PAD, y, CONTENT_W, Inches(0.92),
                 RGBColor(0x1A, 0x14, 0x06), line_color=RGBColor(0x3A, 0x2E, 0x12))

        # Level badge
        lvl_color = RISK_COLORS.get(level, GOLD)
        add_rect(slide, PAD + Inches(0.12), y + Inches(0.26), Inches(0.88), Inches(0.32), lvl_color)
        add_text_box(slide, level, PAD + Inches(0.12), y + Inches(0.26), Inches(0.88), Inches(0.32),
                     font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER, rtl=True, font_name="Inter")

        add_text_box(slide, risk,
                     PAD + Inches(0.18), y + Inches(0.06), CONTENT_W * 0.4, Inches(0.38),
                     font_size=13, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True)

        add_rect(slide, PAD + CONTENT_W * 0.42, y + Inches(0.1),
                 Pt(1.5), Inches(0.72), RGBColor(0x3A, 0x2E, 0x12))

        add_text_box(slide, mitigation,
                     PAD + CONTENT_W * 0.44, y + Inches(0.1),
                     CONTENT_W * 0.54, Inches(0.78),
                     font_size=11, color=GOLD, align=PP_ALIGN.RIGHT, rtl=True, font_name="Inter")


def slide_team(slide, data, slide_num, total):
    slide_bg(slide, BLACK)
    content_header(slide, data["section"], data["title"], slide_num, total)

    # Entity & founder block
    add_rect(slide, PAD, Inches(1.32), CONTENT_W, Inches(1.0),
             RGBColor(0x1A, 0x14, 0x06), line_color=GOLD, line_width=1)
    add_rect(slide, PAD, Inches(1.32), Pt(4), Inches(1.0), GOLD)
    add_text_box(slide, data["entity"],
                 PAD + Inches(0.18), Inches(1.36), CONTENT_W * 0.65, Inches(0.44),
                 font_size=11, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True, font_name="Inter")
    add_text_box(slide, data["founder"],
                 PAD + Inches(0.18), Inches(1.8), CONTENT_W * 0.65, Inches(0.44),
                 font_size=10, color=GOLD, align=PP_ALIGN.RIGHT, rtl=True, font_name="Inter")

    # Team grid
    n = len(data["team"])
    gap = Inches(0.2)
    col_w = (CONTENT_W - gap * (n - 1)) / n
    t_y = Inches(2.44)

    for i, (role, name, detail) in enumerate(data["team"]):
        x = PAD + i * (col_w + gap)
        add_rect(slide, x, t_y, col_w, Inches(2.1),
                 RGBColor(0x18, 0x14, 0x08), line_color=GOLD, line_width=1)
        add_rect(slide, x, t_y, col_w, Inches(0.07), GOLD)
        icon_idx = (i * 7 + 1) % len(ICONS)
        add_image(slide, ICONS[icon_idx], x + col_w / 2 - Inches(0.3), t_y + Inches(0.12),
                  Inches(0.6), Inches(0.6))
        add_text_box(slide, name, x + Inches(0.1), t_y + Inches(0.84), col_w - Inches(0.2), Inches(0.44),
                     font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, rtl=False, font_name="Montserrat")
        add_text_box(slide, role, x + Inches(0.1), t_y + Inches(1.28), col_w - Inches(0.2), Inches(0.36),
                     font_size=10, color=GOLD, align=PP_ALIGN.CENTER, rtl=True, font_name="Inter")
        add_text_box(slide, detail, x + Inches(0.1), t_y + Inches(1.64), col_w - Inches(0.2), Inches(0.4),
                     font_size=9, color=WHITE, align=PP_ALIGN.CENTER, rtl=True, font_name="Inter")

    # Credentials
    cred_y = Inches(4.66)
    add_rect(slide, PAD, cred_y, CONTENT_W, Pt(2), GOLD)
    for i, cred in enumerate(data["credentials"]):
        cy = cred_y + Inches(0.1) + i * Inches(0.42)
        icon_idx = (i * 11 + 5) % len(ICONS)
        add_image(slide, ICONS[icon_idx], CONTENT_R - Inches(0.46), cy + Inches(0.02),
                  Inches(0.38), Inches(0.38))
        add_text_box(slide, cred,
                     PAD, cy, CONTENT_R - Inches(0.6) - PAD, Inches(0.38),
                     font_size=11, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True, font_name="Inter")


def slide_competitors(slide, data, slide_num, total):
    slide_bg(slide, BLACK)
    content_header(slide, data["section"], data["title"], slide_num, total)

    rows = data["comparison"]
    n_cols = len(rows[0])
    n_rows = len(rows)
    col_fracs = [0.28, 0.22, 0.28, 0.22]
    row_h = Inches(0.66)
    start_y = Inches(1.32)
    col_colors = [None, GOLD, RGBColor(0x2A, 0x22, 0x0C), RGBColor(0x1A, 0x14, 0x06)]
    text_colors = [WHITE, BLACK, WHITE, WHITE]

    # Draw table
    for r, row in enumerate(rows):
        y = start_y + r * row_h
        x_cur = PAD
        for c, val in enumerate(row):
            cw = CONTENT_W * col_fracs[c]
            if r == 0:
                add_rect(slide, x_cur, y, cw, row_h, GOLD)
                add_text_box(slide, val, x_cur, y + Inches(0.1), cw, row_h - Inches(0.14),
                             font_size=11, bold=True, color=BLACK,
                             align=PP_ALIGN.CENTER, rtl=True, font_name="Montserrat")
            else:
                bg = col_colors[c] or (RGBColor(0x1A, 0x14, 0x06) if r % 2 == 1 else RGBColor(0x22, 0x1A, 0x08))
                add_rect(slide, x_cur, y, cw, row_h, bg)
                add_text_box(slide, val, x_cur, y + Inches(0.1), cw, row_h - Inches(0.14),
                             font_size=11, color=text_colors[c],
                             align=PP_ALIGN.CENTER, rtl=(c != 3), font_name="Inter")
            x_cur += cw

    # Conclusion
    conc_y = start_y + n_rows * row_h + Inches(0.1)
    add_rect(slide, PAD, conc_y, CONTENT_W, Inches(0.56),
             RGBColor(0x14, 0x10, 0x04), line_color=GOLD, line_width=1)
    add_text_box(slide, data["conclusion"],
                 PAD + Inches(0.15), conc_y + Inches(0.08), CONTENT_W - Inches(0.3), Inches(0.42),
                 font_size=11, color=GOLD, align=PP_ALIGN.RIGHT, rtl=True, font_name="Inter")


# ═══════════════════════════════════════════════════════════════════════════════
#  MAIN PPTX BUILD
# ═══════════════════════════════════════════════════════════════════════════════

def build_pptx(output_path):
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]
    total = len(SLIDES)

    builders = {
        "cover":            slide_cover,
        "toc":              slide_toc,
        "exec_summary":     slide_exec_summary,
        "section_divider":  slide_section_divider,
        "problem":          slide_problem,
        "solution":         slide_solution,
        "curriculum":       slide_curriculum,
        "proof":            slide_proof,
        "employers":        slide_employers,
        "phases":           slide_phases,
        "budget":           slide_budget,
        "kpis":             slide_kpis,
        "mne":              slide_mne,
        "risk":             slide_risk,
        "team":             slide_team,
        "competitors":      slide_competitors,
        "alignment":        slide_alignment,
        "closing":          slide_closing,
    }

    for idx, data in enumerate(SLIDES, 1):
        slide = prs.slides.add_slide(blank)
        stype = data["type"]
        fn = builders.get(stype)
        if fn:
            if stype == "cover":
                fn(slide, data)
            else:
                fn(slide, data, idx, total)
        print(f"  [{idx:02d}/{total}] {stype}")

    prs.save(output_path)
    print(f"\nSaved: {output_path}")


# ═══════════════════════════════════════════════════════════════════════════════
#  DOCX BUILD
# ═══════════════════════════════════════════════════════════════════════════════

def set_rtl_doc(para):
    pPr = para._p.get_or_add_pPr()
    bidi = OxmlElement('w:bidi')
    bidi.set(dqn('w:val'), '1')
    pPr.append(bidi)
    jc = OxmlElement('w:jc')
    jc.set(dqn('w:val'), 'right')
    pPr.append(jc)


def doc_heading(doc, text, level=1, color=DGOLD):
    p = doc.add_paragraph()
    set_rtl_doc(p)
    p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    run = p.add_run(text)
    run.font.size = DPt(26 - (level - 1) * 4)
    run.font.bold = True
    run.font.color.rgb = color
    run.font.name = "Montserrat"
    return p


def doc_para(doc, text, color=DWHITE, size=11):
    p = doc.add_paragraph()
    set_rtl_doc(p)
    p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    run = p.add_run(text)
    run.font.size = DPt(size)
    run.font.color.rgb = color
    run.font.name = "Arial"
    return p


def doc_bullet(doc, text):
    p = doc.add_paragraph(style='List Bullet')
    set_rtl_doc(p)
    p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    run = p.add_run(text)
    run.font.size = DPt(11)
    run.font.color.rgb = DBLACK
    run.font.name = "Arial"
    return p


def build_docx(output_path):
    doc = Document()
    # Page background approximation via page color
    section = doc.sections[0]
    section.page_width = DInches(8.5)
    section.page_height = DInches(11)
    section.left_margin = DInches(1.0)
    section.right_margin = DInches(1.0)
    section.top_margin = DInches(0.9)
    section.bottom_margin = DInches(0.9)

    # Title page
    doc_heading(doc, "MY4 Education", level=1, color=DBLACK)
    doc_heading(doc, "مقترح شراكة استراتيجية", level=1, color=DBLACK)
    doc_para(doc, "وزارة التضامن الاجتماعي — جمهورية مصر العربية", color=DRGBColor(0xC8, 0xA2, 0x4C))
    doc_para(doc, "من قاعة الدارس إلى سوق العمل  ·  ٢٠٢٦", color=DRGBColor(0x44, 0x44, 0x44), size=10)
    doc.add_page_break()

    for data in SLIDES:
        stype = data["type"]

        if stype == "section_divider":
            doc.add_page_break()
            doc_heading(doc, f"{data['num']}  ·  {data['title']}", level=1, color=DBLACK)
            doc_para(doc, data["subtitle"], color=DRGBColor(0x66, 0x55, 0x22), size=12)
            doc.add_paragraph("─" * 60)

        elif stype == "exec_summary":
            doc_heading(doc, data["title"], level=2, color=DBLACK)
            doc_para(doc, data["body"])
            doc.add_paragraph()
            for num, label, src in data["stats"]:
                doc_para(doc, f"{num}  —  {label}  ({src})",
                         color=DRGBColor(0xC8, 0xA2, 0x4C), size=12)

        elif stype == "problem":
            doc_heading(doc, data["title"], level=2, color=DBLACK)
            for num, label, src in data["stats"]:
                doc_bullet(doc, f"{num}  {label}  — {src}")

        elif stype == "solution":
            doc_heading(doc, data["title"], level=2, color=DBLACK)
            for num, title, body in data["steps"]:
                doc_heading(doc, f"{num}  {title}", level=3, color=DRGBColor(0x88, 0x66, 0x22))
                doc_para(doc, body)
            doc.add_paragraph()
            doc_para(doc, data["why"], color=DRGBColor(0x33, 0x33, 0x33), size=10)

        elif stype == "curriculum":
            doc_heading(doc, data["title"], level=2, color=DBLACK)
            for period, topic, detail in data["weeks"]:
                doc_heading(doc, f"{period} — {topic}", level=3, color=DRGBColor(0x88, 0x66, 0x22))
                doc_para(doc, detail)
            doc.add_paragraph()
            doc_para(doc, data["female_track"], color=DRGBColor(0xA4, 0x23, 0x2A), size=11)
            doc_para(doc, data["cert"], color=DRGBColor(0x44, 0x44, 0x44), size=10)

        elif stype == "proof":
            doc_heading(doc, data["title"], level=2, color=DBLACK)
            doc_para(doc, data["body"])
            for b in data["bullets"]:
                doc_bullet(doc, b)

        elif stype == "employers":
            doc_heading(doc, data["title"], level=2, color=DBLACK)
            for sector, companies in data["sectors"]:
                doc_heading(doc, sector, level=3, color=DRGBColor(0x88, 0x66, 0x22))
                doc_para(doc, " · ".join(companies))
            doc_para(doc, data["note"], color=DRGBColor(0x66, 0x55, 0x22), size=10)

        elif stype == "phases":
            doc_heading(doc, data["title"], level=2, color=DBLACK)
            for phase in data["phases"]:
                doc_heading(doc, f"{phase['name']}  ({phase['duration']} — {phase['count']})",
                            level=3, color=DRGBColor(0x88, 0x66, 0x22))
                for pt in phase["points"]:
                    doc_bullet(doc, pt)

        elif stype == "budget":
            doc_heading(doc, data["title"], level=2, color=DBLACK)
            doc_para(doc, data["subtitle"], color=DRGBColor(0xC8, 0xA2, 0x4C), size=11)
            for item, amount, pct in data["items"]:
                doc_bullet(doc, f"{item}:  {amount} جنيه  ({pct})")
            doc_para(doc, f"الإجمالي:  {data['total']} جنيه  |  {data['per_trainee']} جنيه / متدرب",
                     color=DRGBColor(0xC8, 0xA2, 0x4C), size=13)
            doc_para(doc, data["roi"])
            doc_para(doc, data["clawback"], color=DRGBColor(0xA4, 0x23, 0x2A), size=11)

        elif stype == "kpis":
            doc_heading(doc, data["title"], level=2, color=DBLACK)
            for num, label, sub in data["kpis"]:
                doc_para(doc, f"{num}  —  {label}:  {sub}",
                         color=DRGBColor(0xC8, 0xA2, 0x4C), size=12)
            doc.add_paragraph()
            doc_heading(doc, "التزامات إضافية", level=3, color=DBLACK)
            for q in data["qualitative"]:
                doc_bullet(doc, q)

        elif stype == "mne":
            doc_heading(doc, data["title"], level=2, color=DBLACK)
            for cat, rows in data["indicators"]:
                doc_heading(doc, cat, level=3, color=DRGBColor(0x88, 0x66, 0x22))
                for indicator, target, method in rows:
                    doc_bullet(doc, f"{indicator}  |  {target}  |  {method}")
            doc_para(doc, data["reporting"], color=DRGBColor(0x44, 0x44, 0x44), size=10)

        elif stype == "risk":
            doc_heading(doc, data["title"], level=2, color=DBLACK)
            for risk, level, mitigation in data["risks"]:
                doc_heading(doc, f"[{level}]  {risk}", level=3, color=DRGBColor(0x88, 0x66, 0x22))
                doc_para(doc, mitigation)

        elif stype == "team":
            doc_heading(doc, data["title"], level=2, color=DBLACK)
            doc_para(doc, data["entity"])
            doc_para(doc, data["founder"], color=DRGBColor(0xC8, 0xA2, 0x4C), size=11)
            doc.add_paragraph()
            for role, name, detail in data["team"]:
                doc_bullet(doc, f"{name}  —  {role}:  {detail}")
            doc.add_paragraph()
            doc_heading(doc, "الاعتمادات", level=3, color=DBLACK)
            for cred in data["credentials"]:
                doc_bullet(doc, cred)

        elif stype == "competitors":
            doc_heading(doc, data["title"], level=2, color=DBLACK)
            for row in data["comparison"]:
                doc_para(doc, "  |  ".join(row))
            doc.add_paragraph()
            doc_para(doc, data["conclusion"], color=DRGBColor(0xC8, 0xA2, 0x4C), size=11)

        elif stype == "alignment":
            doc_heading(doc, data["title"], level=2, color=DBLACK)
            for prog, desc in data["alignment"]:
                doc_para(doc, f"▸  {prog}:  {desc}", color=DRGBColor(0x44, 0x44, 0x44))
            doc.add_paragraph()
            doc_heading(doc, "المطلوب من الوزارة", level=3, color=DRGBColor(0xA4, 0x23, 0x2A))
            for req in data["requests"]:
                doc_heading(doc, f"{req['num']} — {req['title']}", level=3,
                            color=DRGBColor(0xA4, 0x23, 0x2A))
                doc_para(doc, req["body"])

        elif stype == "closing":
            doc.add_page_break()
            doc_heading(doc, data["title"], level=1, color=DBLACK)
            doc_para(doc, data["tagline"], color=DRGBColor(0xC8, 0xA2, 0x4C), size=14)
            doc.add_paragraph()
            for label, val in data["contact"]:
                doc_para(doc, f"{label}:  {val}", color=DBLACK, size=11)

    doc.save(output_path)
    print(f"Saved: {output_path}")


if __name__ == "__main__":
    out = "/home/user/Mahmoud-Gadalla/outputs"
    os.makedirs(out, exist_ok=True)
    print("Building PPTX...")
    build_pptx(f"{out}/MY4_Education_Ministry_Proposal.pptx")
    print("Building DOCX...")
    build_docx(f"{out}/MY4_Education_Ministry_Proposal.docx")
