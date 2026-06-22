import copy, re
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

SRC  = "/root/.claude/uploads/1c498b9e-7d15-5f41-80e6-6ba18c7c672b/5b244a6a-MY4_Education_Brand_Guidelines.pptx"
DEST = "/home/user/Mahmoud-Gadalla/MY4_Education_Ministry_Proposal.pptx"

# ── helpers ──────────────────────────────────────────────────────────────────

def copy_slide(prs, slide_index):
    """Duplicate a slide from the same presentation and append it."""
    template_slide = prs.slides[slide_index]
    slide_layout   = template_slide.slide_layout
    new_slide      = prs.slides.add_slide(slide_layout)

    # Copy all XML children from template body
    sp_tree = new_slide.shapes._spTree
    for el in list(sp_tree):
        sp_tree.remove(el)
    for el in template_slide.shapes._spTree:
        sp_tree.append(copy.deepcopy(el))

    # Copy background
    bg = template_slide._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}bg')
    if bg is not None:
        existing = new_slide._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}bg')
        if existing is None:
            cSld = new_slide._element.find('{http://schemas.openxmlformats.org/presentationml/2006/main}cSld')
            if cSld is not None:
                cSld.insert(0, copy.deepcopy(bg))
    return new_slide


def replace_text(slide, replacements):
    """Replace text in all shapes on a slide. replacements = {old: new}"""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                for old, new in replacements.items():
                    if old in run.text:
                        run.text = run.text.replace(old, new)
            # Also handle full paragraph text spanning single run
            full = "".join(r.text for r in para.runs)
            for old, new in replacements.items():
                if old in full and len(para.runs) == 1:
                    para.runs[0].text = full.replace(old, new)


def set_slide_text(slide, mapping):
    """mapping: {exact_current_text: new_text}  — replaces whole shape text."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        full = shape.text_frame.text.strip()
        if full in mapping:
            tf = shape.text_frame
            # preserve formatting of first run
            for para in tf.paragraphs:
                for run in para.runs:
                    run.text = ""
            # set text in first paragraph first run
            if tf.paragraphs and tf.paragraphs[0].runs:
                tf.paragraphs[0].runs[0].text = mapping[full]
            else:
                tf.paragraphs[0].text = mapping[full]


# ── load source ───────────────────────────────────────────────────────────────
prs = Presentation(SRC)
total_src = len(prs.slides)

# We will build a fresh presentation by loading the source and removing
# all existing slides, then re-adding only what we need.
# python-pptx doesn't allow removing slides cleanly, so we'll use a different
# approach: copy selected slides into a new file via XML surgery.

from pptx.opc.constants import RELATIONSHIP_TYPE as RT
import zipfile, shutil, os, tempfile

# ────────────────────────────────────────────────────────────────────────────
# Build proposal by copying slide XMLs from source into a new file
# ────────────────────────────────────────────────────────────────────────────

# Slides to copy (0-indexed from source):
# 0  = Cover (dark full-bleed)
# 9  = TOC (table of contents)
# 10 = Section divider
# 11 = 2-col goal/aim  (light bg)
# 12 = 3-col resources (mixed)
# 13 = Stats hero (dark)
# 14 = 5-step process (light)
# 15 = Table (light)
# 16 = Chart (dark)
# 17 = Thank you (dark)

SLIDE_MAP = [0, 9, 10, 13, 10, 11, 12, 10, 14, 11, 10, 17]
# 0=Cover, 9=TOC, 10=Section(Problem), 13=Stats, 10=Section(Solution),
# 11=2col(solution detail), 12=3col(model steps), 10=Section(Roadmap),
# 14=5step(phases), 11=2col(ask+KPIs), 10=Section(Contact), 17=Thankyou

print("Copying slides...")
new_slides = []
for idx in SLIDE_MAP:
    s = copy_slide(prs, idx)
    new_slides.append(s)

# Remove original 18 slides (they're now at the front)
# We keep only the newly added ones
rId_list = prs.slides._sldIdLst
slide_elements = list(rId_list)
# Newly added slides are the last len(SLIDE_MAP) entries
to_keep = slide_elements[total_src:]
to_remove = slide_elements[:total_src]

for el in to_remove:
    rId_list.remove(el)

print(f"Kept {len(list(rId_list))} slides")

# ── Now update slide content ──────────────────────────────────────────────────
slides = list(prs.slides)

BRAND_FOOTER = {"MERIDIAN   BRAND SYSTEM": "MY4 EDUCATION",
                "MERIDIAN  BRAND SYSTEM":  "MY4 EDUCATION",
                "MERIDIAN BRAND SYSTEM":   "MY4 EDUCATION"}

# Helper to apply brand footer fix on all slides
def fix_footer(slide, num, total=12):
    replace_text(slide, BRAND_FOOTER)
    replace_text(slide, {
        "01 / 18": f"01 / {total}",
        "02 / 18": f"02 / {total}",
        "03 / 18": f"03 / {total}",
        "04 / 18": f"04 / {total}",
        "05 / 18": f"05 / {total}",
        "06 / 18": f"06 / {total}",
        "07 / 18": f"07 / {total}",
        "08 / 18": f"08 / {total}",
        "09 / 18": f"09 / {total}",
        "10 / 18": f"10 / {total}",
        "11 / 18": f"11 / {total}",
        "12 / 18": f"12 / {total}",
        "13 / 18": f"13 / {total}",
        "14 / 18": f"14 / {total}",
        "15 / 18": f"15 / {total}",
        "16 / 18": f"16 / {total}",
        "17 / 18": f"17 / {total}",
        "18 / 18": f"18 / {total}",
    })
    replace_text(slide, {f"XX / {total}": f"{num:02d} / {total}"})

TOTAL = len(slides)

# ── Slide 1: Cover ───────────────────────────────────────────────────────────
s = slides[0]
replace_text(s, {
    "MERIDIAN": "MY4 EDUCATION",
    "Brand Visual\nGuidelines &\nPresentation System": "مبادرة وطنية\nلتأهيل الشباب الجامعي\nوتمكينهم اقتصادياً",
    "Brand Visual": "مبادرة وطنية",
    "Guidelines &": "لتأهيل الشباب الجامعي",
    "Presentation System": "وتمكينهم اقتصادياً",
    "The complete identity reference — palette, typography, the woven mark and a ready-to-use slide template. Built on a Black · White · Gold · Red system.": 
        "مقدَّم إلى: وزارة التضامن الاجتماعي | يونيو ٢٠٢٦",
    "01 — 09  IDENTITY        10 — 18  TEMPLATE": "Gadalla111@gmail.com  |  +20 111 037 3331",
})
replace_text(s, BRAND_FOOTER)
replace_text(s, {"01 / 18": "01 / 12", "M\nMERIDIAN": "MY4\nEDUCATION"})

# ── Slide 2: Table of Contents ────────────────────────────────────────────────
s = slides[1]
replace_text(s, {
    "PRESENTATION TEMPLATE": "جدول المحتويات",
    "Table of contents": "مقترح MY4 Education — وزارة التضامن الاجتماعي",
    "Objectives": "الإشكالية",
    "Goals, aim and the outcome we are working toward.": "الشهادة وحدها لا تصنع فرصة عمل — الأرقام تُحدِّد الأزمة",
    "Resources": "الحل",
    "Team, budget and the assets the work depends on.": "تدريب حقيقي · إثبات مُصوَّر",
    "Budget": "إثبات النموذج",
    "Where funds come from and where they are spent.": "تجربة ريادية فعلية — الأكاديمية العربية (AAST)",
    "Solutions": "خطة المراحل",
    "The proposed approach and why it is the right one.": "من ٤٥ شاباً إلى برنامج وطني",
    "Comms. plan": "المؤشرات والتوافق",
    "How progress is shared with stakeholders.": "أرقام نُقدِّمها ونتحمل مسؤوليتها",
    "Conclusion": "الطلب من الوزارة",
    "The decision we are asking for, and next steps.": "رعاية رسمية · تمويل تجريبي · وصول تشغيلي",
})
fix_footer(s, 2)

# ── Slide 3: Section — Problem ───────────────────────────────────────────────
s = slides[2]
replace_text(s, {
    "SECTION": "القسم الأول",
    "01": "01",
    "Objectives": "الإشكالية",
    "What the project sets out to achieve, and how we'll know it worked.": "الشهادة وحدها لا تصنع فرصة عمل",
})
fix_footer(s, 3)

# ── Slide 4: Stats ────────────────────────────────────────────────────────────
s = slides[3]
replace_text(s, {
    "03 · BUDGET": "الأرقام تُحدِّد الأزمة بدقة",
    "The numbers that matter": "مصادر: CAPMAS 2025 · Nexford Egypt 2021",
    "$10M": "6.3%",
    "Projected first-year revenue": "معدل البطالة الإجمالي — مصر ٢٠٢٥",
    "Based on a 200-unit break-even and a 25% target market share.": "المصدر: CAPMAS، الربع الرابع ٢٠٢٥",
    "$500K": "45%",
    "Net profit of the project": "من إجمالي العاطلين خريجو جامعات",
    "200": "78%",
    "Units to break even": "من الشركات تعجز عن إيجاد الكفاءات",
    "25%": "70%",
    "Target market share": "من أصحاب العمل سيوظفون أكثر لو كان الخريجون جاهزين",
})
fix_footer(s, 4)

# ── Slide 5: Section — Solution ──────────────────────────────────────────────
s = slides[4]
replace_text(s, {
    "SECTION": "القسم الثاني",
    "01": "02",
    "Objectives": "الحل",
    "What the project sets out to achieve, and how we'll know it worked.": "تدريب حقيقي · إثبات مُصوَّر",
})
fix_footer(s, 5)

# ── Slide 6: 2-col — Solution detail ─────────────────────────────────────────
s = slides[5]
replace_text(s, {
    "01 · OBJECTIVES": "كيف يعمل النموذج",
    "The goal & the aim": "ثلاث خطوات تُحوِّل الخريج إلى محترف جاهز",
    "The goal": "التدريب الميداني",
    "Goals are specific and measurable, with clear deadlines and outcomes. They focus the organisation and make sure resources are spent where they matter most.":
        "شهر كامل داخل شركة مضيفة حقيقية تحت إشراف متخصص — منهج منظَّم ومهام واقعية تبني مهارات قابلة للتطبيق الفوري.",
    "Our aim": "الإثبات المُصوَّر",
    "The aim is the result the organisation seeks to achieve — clear, achievable and the basis for the planning that follows. It is the north star for the work ahead.":
        "امتحان عملي ميداني مصوَّر بالكامل يُوثِّق قدرة كل مشارك على الأداء — دليل موضوعي قابل للقياس يُرضي أصحاب العمل والوزارة.",
})
fix_footer(s, 6)

# ── Slide 7: 3-col — Model steps / Proof ─────────────────────────────────────
s = slides[6]
replace_text(s, {
    "02 · RESOURCES": "إثبات النموذج — الأكاديمية العربية",
    "What the work depends on": "نموذج مُجرَّب لا نظرية على الورق",
    "Human resources": "التجربة الريادية",
    "An experienced team with the skills to deliver on time and on budget. Clear ownership at every stage.":
        "تجربة فعلية داخل AAST في تخصص إدارة المخازن واللوجستيات — الكليات المستهدفة: اللوجستيات · التسويق · الإعلام.",
    "Financial resources": "المرحلة التجريبية",
    "A working budget of $100,000 covering salaries, equipment and contingency for the full scope.":
        "جامعة واحدة · ٣ كليات · ٤٥ شاباً · دورة واحدة · شركتان مضيفتان · شهر تدريب + امتحان ختامي مصوَّر.",
    "Physical resources": "الأثر الدولي",
    "Specialist equipment, tested and calibrated before use, plus the systems the team relies on daily.":
        "وزارة العمل الأمريكية: 93% من مكتملي برامج التلمذة يحتفظون بوظائفهم. NC Commerce 2026: معدل توظيف 90% مقابل 74% — فارق 16 نقطة يستمر عقداً.",
})
fix_footer(s, 7)

# ── Slide 8: Section — Roadmap ───────────────────────────────────────────────
s = slides[7]
replace_text(s, {
    "SECTION": "القسم الثالث",
    "01": "03",
    "Objectives": "خطة المراحل",
    "What the project sets out to achieve, and how we'll know it worked.": "من ٤٥ شاباً إلى برنامج وطني",
})
fix_footer(s, 8)

# ── Slide 9: 5-step — Phases ─────────────────────────────────────────────────
s = slides[8]
replace_text(s, {
    "04 · SOLUTIONS": "خطة المراحل الثلاث",
    "How we get there": "من التجريب إلى التوسع الوطني",
    "1": "01",
    "Research": "التجريب",
    "Study the market and audience to find the opening.": "٠–٦ أشهر · جامعة واحدة · ٣ كليات · ٤٥ شاباً",
    "2": "02",
    "Build": "التوسع",
    "Define features, design and price; create the product.": "٦–١٨ شهراً · ٤ جامعات · ~٥٠٠ شاب",
    "3": "03",
    "Test": "الوطني",
    "Trial with real users and refine on their feedback.": "١٨–٣٦ شهراً · +١٥ جامعة · +٣٠٠٠ شاب سنوياً",
    "4": "04",
    "Launch": "القياس",
    "Release to market with a full go-to-market plan.": "تقارير أثر نصف سنوية · بيانات توظيف موثَّقة",
    "5": "05",
    "Sustain": "الاستدامة",
    "Support customers and improve on what we learn.": "دمج مع وحدات التضامن الجامعية في ٤٣ جامعة",
})
fix_footer(s, 9)

# ── Slide 10: 2-col — KPIs + Alignment ───────────────────────────────────────
s = slides[9]
replace_text(s, {
    "01 · OBJECTIVES": "المؤشرات والتوافق مع الوزارة",
    "The goal & the aim": "أرقام نُقدِّمها ونتحمل مسؤوليتها",
    "The goal": "آلية القياس",
    "Goals are specific and measurable, with clear deadlines and outcomes. They focus the organisation and make sure resources are spent where they matter most.":
        "• تقييم جاهزية التوظيف قبل وبعد\n• تتبُّع عروض العمل خلال ٦ أشهر\n• استطلاع رضا جهات التشغيل\n• تقرير أثر نصف سنوي للوزارة",
    "Our aim": "التوافق مع الأولويات",
    "The aim is the result the organisation seeks to achieve — clear, achievable and the basis for the planning that follows. It is the north star for the work ahead.":
        "✓ وحدات التضامن الجامعية · ٤٣ جامعة\n✓ رؤية مصر ٢٠٣٠\n✓ الإطار المرجعي للتعليم العالي SCU 2025\n✓ قمة START 2026 — وزارة التضامن",
})
fix_footer(s, 10)

# ── Slide 11: Section — Ask ──────────────────────────────────────────────────
s = slides[10]
replace_text(s, {
    "SECTION": "القسم الرابع",
    "01": "04",
    "Objectives": "الطلب من الوزارة",
    "What the project sets out to achieve, and how we'll know it worked.": "رعاية رسمية · تمويل تجريبي · وصول تشغيلي",
})
fix_footer(s, 11)

# ── Slide 12: Thank you / Contact ────────────────────────────────────────────
s = slides[11]
replace_text(s, {
    "Thank you.": "شكراً.",
    "MERIDIAN is yours to build on — duplicate any layout, swap the wordmark, keep the system.":
        "MY4 Education — مبادرة وطنية لتأهيل الشباب الجامعي وتمكينهم اقتصادياً\nأولاً: الرعاية الرسمية · ثانياً: التمويل التجريبي · ثالثاً: الوصول التشغيلي",
    "MERIDIAN": "MY4 Education",
    "EMAIL": "البريد الإلكتروني",
    "hello@meridian.brand": "Gadalla111@gmail.com",
    "WEB": "المؤسس",
    "meridian.brand": "Mahmoud Gadalla",
    "PHONE": "الهاتف",
    "+20 100 000 0000": "+20 111 037 3331",
    "18 / 18": "12 / 12",
})
replace_text(s, BRAND_FOOTER)

print("Saving...")
prs.save(DEST)
print(f"Saved → {DEST}")
